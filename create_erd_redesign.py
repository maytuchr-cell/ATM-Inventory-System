"""
ATM Inventory System — ERD Redesign (v2)
แยก "ตัวอะไหล่ (Master)" ออกจาก "อะไหล่ในคลัง (Stock)" ให้ทีมเข้าใจง่าย
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

NAVY   = RGBColor(0x1C, 0x35, 0x57)
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # Master
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
ORANGED= RGBColor(0xD8, 0x5A, 0x30)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
TEXT   = RGBColor(0x1E, 0x2D, 0x40)
TEAL   = RGBColor(0x0D, 0x94, 0x88)   # Inventory
GREEN  = RGBColor(0x05, 0x96, 0x69)
RED    = RGBColor(0xDC, 0x26, 0x26)
AMBER  = RGBColor(0xB4, 0x53, 0x09)   # Serial
GREY   = RGBColor(0x6B, 0x72, 0x80)   # Ledger
PURPLE = RGBColor(0x53, 0x4A, 0xB7)
GREY_L = RGBColor(0xD1, 0xD5, 0xDB)
ROW_ALT= RGBColor(0xF3, 0xF6, 0xFB)
REDBG  = RGBColor(0xFD, 0xEC, 0xEC)
GRNBG  = RGBColor(0xEA, 0xF6, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])
def bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def box(slide, x, y, w, h, fill, line=None, lw=0.75, radius=False):
    sh = slide.shapes.add_shape(5 if radius else 1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(slide, text, x, y, w, h, size=12, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, vmid=False, font="Consolas"):
    if color is None: color = TEXT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.03); tf.margin_right=Inches(0.03); tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01)
    if vmid: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def header(slide, title, subtitle=None):
    box(slide, 0, 0, 10, 0.78, NAVY)
    box(slide, 0, 0.78, 10, 0.04, ORANGE)
    txt(slide, title, 0.4, 0.08, 9.2, 0.55, size=20, bold=True, color=WHITE, vmid=True, font="Calibri")
    if subtitle:
        txt(slide, subtitle, 0.42, 0.88, 9, 0.26, size=10, color=MUTED, italic=True, font="Calibri")

def field_table(slide, x, y, fields, col_w=(1.7, 1.55, 0.5, 0.5, 4.45), row_h=0.30, fsize=8.5, header_clr=NAVY):
    cols = ["Field", "Data Type", "Null", "Key", "Description"]
    tw = sum(col_w)
    box(slide, x, y, tw, 0.3, header_clr)
    cx = x
    for i, c in enumerate(cols):
        al = PP_ALIGN.CENTER if c in ("Null","Key") else PP_ALIGN.LEFT
        txt(slide, c, cx+0.05, y+0.02, col_w[i]-0.08, 0.26, size=8.5, bold=True, color=WHITE, vmid=True, font="Calibri", align=al)
        cx += col_w[i]
    box(slide, x, y+0.3, tw, len(fields)*row_h, WHITE, GREY_L, 0.75)
    for ri, (fname, ftype, nullable, key, desc) in enumerate(fields):
        ry = y + 0.3 + ri*row_h
        if ri % 2 == 1: box(slide, x+0.02, ry, tw-0.04, row_h, ROW_ALT)
        cx = x
        txt(slide, fname, cx+0.06, ry, col_w[0]-0.1, row_h, size=fsize, bold=True, color=TEXT, vmid=True)
        cx += col_w[0]
        txt(slide, ftype, cx+0.06, ry, col_w[1]-0.1, row_h, size=fsize-0.5, color=BLUE, vmid=True)
        cx += col_w[1]
        nl = "NULL" if nullable else "NOT"
        nc = MUTED if nullable else RGBColor(0x9a,0x9a,0x9a)
        txt(slide, nl, cx+0.02, ry, col_w[2]-0.04, row_h, size=7, color=nc, vmid=True, align=PP_ALIGN.CENTER, font="Calibri")
        cx += col_w[2]
        kc = {"PK":ORANGED, "UK":PURPLE, "FK":GREEN}.get(key, MUTED)
        if key: txt(slide, key, cx+0.03, ry, col_w[3]-0.06, row_h, size=fsize-0.5, bold=True, color=kc, vmid=True, align=PP_ALIGN.CENTER, font="Calibri")
        cx += col_w[3]
        txt(slide, desc, cx+0.06, ry, col_w[4]-0.1, row_h, size=fsize, color=TEXT, vmid=True, font="Calibri")

def legend(slide, y=5.32):
    items = [("PK", ORANGED, "Primary"), ("UK", PURPLE, "Unique"), ("FK", GREEN, "Foreign")]
    x = 0.4
    for tag, clr, desc in items:
        txt(slide, tag, x, y-0.02, 0.4, 0.25, size=9, bold=True, color=clr, vmid=True, font="Calibri"); x += 0.34
        txt(slide, desc, x, y-0.02, 1.2, 0.25, size=9, color=MUTED, vmid=True, font="Calibri"); x += len(desc)*0.06 + 0.25
    txt(slide, "Null = nullable column", x, y-0.02, 2.2, 0.25, size=9, color=MUTED, vmid=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, NAVY)
box(s, 0, 0, 0.55, 5.625, ORANGE)
txt(s, "Database Redesign", 0.95, 0.95, 8.5, 1.0, size=38, bold=True, color=WHITE, font="Calibri")
txt(s, "แยก \"ตัวอะไหล่\" ออกจาก \"อะไหล่ในคลัง\"", 0.97, 2.0, 8.5, 0.5, size=18, color=RGBColor(0xB0,0xC8,0xE8), font="Calibri")
box(s, 0.97, 2.75, 7.3, 0.045, ORANGE)
txt(s, "Part (Master)  ·  PartStock (Balance)  ·  PartUnit (Serial)  ·  StockMovement (Ledger)",
    0.97, 2.9, 8.4, 0.4, size=13, color=RGBColor(0x90,0xAE,0xD8), font="Calibri")
txt(s, "ATM Inventory System  |  DataOne Asia  |  June 2026", 0.97, 3.62, 8, 0.36, size=12, color=MUTED, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — สารบัญ (Table of Contents)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "สารบัญ", "Table of Contents")

toc = [
    ("01", "System Architecture — ภาพรวมทั้งระบบ", "Browser · Frontend · API · Database", NAVY),
    ("02", "ปัญหา — ทำไมตาราง Part เดิมทำให้ทีมงง", "ปนกัน 3 ความหมายในตารางเดียว", RED),
    ("03", "แนวคิดใหม่ — แยกเป็น 3 ชั้น + 1 ประวัติ", "แต่ละตารางมีหน้าที่เดียว ชัดเจน", BLUE),
    ("04", "ERD — แผนภาพความสัมพันธ์ทั้งหมด", "กล่อง = ตาราง · เส้น = ความสัมพันธ์", PURPLE),
    ("05", "Part — ตัวอะไหล่ (Master / Catalog)", "นิยามของอะไหล่ — คืออะไร", BLUE),
    ("06", "PartStock — จำนวนในคลัง (Balance)", "ยอดคงเหลือแยกคลัง ดี/เสีย", TEAL),
    ("07", "PartUnit — ชิ้นเฉพาะตัว (Serial)", "track รายชิ้น (ทางเลือก)", AMBER),
    ("08", "StockMovement — ประวัติการเคลื่อนไหว", "Ledger ทุกการขยับของ", GREY),
    ("09", "Before → After — ฟิลด์ไหนย้ายไปไหน", "สรุปการเปลี่ยนแปลง", GREEN),
    ("10", "ความสัมพันธ์ & แผนย้ายข้อมูล", "3 ขั้นแบบไม่เสียข้อมูล", ORANGED),
]
col_x = [0.45, 5.05]
yy0 = 1.15
for i, (num, title, sub, clr) in enumerate(toc):
    col = i // 5
    row = i % 5
    x = col_x[col]
    y = yy0 + row*0.80
    box(s, x, y, 4.5, 0.68, WHITE, GREY_L, 0.75, radius=True)
    box(s, x, y, 0.08, 0.68, clr)
    box(s, x+0.18, y+0.13, 0.46, 0.42, clr, radius=True)
    txt(s, num, x+0.18, y+0.13, 0.46, 0.42, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, vmid=True, font="Calibri")
    txt(s, title, x+0.78, y+0.08, 3.6, 0.3, size=10.5, bold=True, color=TEXT, vmid=True, font="Calibri")
    txt(s, sub, x+0.78, y+0.37, 3.6, 0.26, size=8.5, color=MUTED, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "System Architecture — ภาพรวมทั้งระบบ")

def darrow(slide, x, y, w=0.32, h=0.34, clr=MUTED):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def chip(slide, x, y, w, h, label, clr, bgc, fs=8.5):
    box(slide, x, y, w, h, bgc, clr, 0.75, radius=True)
    txt(slide, label, x+0.05, y, w-0.1, h, size=fs, bold=True, color=clr, align=PP_ALIGN.CENTER, vmid=True, font="Calibri")

# ── Tier 1: Users ──
box(s, 0.3, 0.98, 9.4, 0.6, RGBColor(0xEE,0xF1,0xF6), GREY_L, 0.75, radius=True)
txt(s, "ผู้ใช้งาน", 0.42, 1.02, 1.5, 0.5, size=9, bold=True, color=MUTED, vmid=True, font="Calibri")
chip(s, 2.7, 1.08, 2.1, 0.4, "👤  Admin (Browser)", NAVY, WHITE, 9)
chip(s, 5.2, 1.08, 2.1, 0.4, "🔧  Technician (Browser)", NAVY, WHITE, 9)

darrow(s, 4.85, 1.62, clr=GREEN)
txt(s, "เปิดหน้าเว็บ (HTTP)", 5.25, 1.66, 2.2, 0.25, size=8, color=GREEN, italic=True, vmid=True, font="Calibri")

# ── Tier 2: Frontend ──
box(s, 0.3, 2.02, 9.4, 0.66, RGBColor(0xEE,0xF3,0xFE), BLUE, 1.0, radius=True)
txt(s, "Frontend  —  Static Web  (IIS ภายในองค์กร / dev: http.server :3000)", 0.45, 2.06, 6.5, 0.3, size=9.5, bold=True, color=BLUE, font="Calibri")
fchips = ["login.html", "admin*.html", "tech.html", "shared/api.js", "translations.js", "styles.css · DM Sans"]
cx = 0.45
for c in fchips:
    w = 0.18 + len(c)*0.072
    chip(s, cx, 2.36, w, 0.26, c, BLUE, WHITE, 8)
    cx += w + 0.12

darrow(s, 4.85, 2.74, clr=ORANGED)
txt(s, "fetch + JWT Bearer token  (Authorization header)", 5.25, 2.78, 4.0, 0.25, size=8, color=ORANGED, italic=True, vmid=True, font="Calibri")

# ── Tier 3: Backend API ──
box(s, 0.3, 3.16, 9.4, 1.18, WHITE, NAVY, 1.0, radius=True)
txt(s, "Backend API  —  ASP.NET Core (.NET 10)  ·  http://...:5128", 0.45, 3.2, 7.0, 0.3, size=9.5, bold=True, color=NAVY, font="Calibri")
# Controllers row
txt(s, "Controllers", 0.45, 3.52, 1.2, 0.22, size=8, bold=True, color=MUTED, font="Calibri")
ctrls = ["Auth", "Parts", "Categories", "GoodsReceipt", "Tickets", "Returns", "Transfers", "StockCount", "Disposal", "Dashboard", "Reports"]
cx = 1.55
for c in ctrls:
    w = 0.16 + len(c)*0.066
    chip(s, cx, 3.50, w, 0.24, c, GREEN, RGBColor(0xEC,0xF7,0xF1), 7.5)
    cx += w + 0.08
# Services + EF row
txt(s, "Services", 0.45, 3.84, 1.2, 0.22, size=8, bold=True, color=MUTED, font="Calibri")
for i, c in enumerate(["JwtHelper", "PasswordHasher", "StockService"]):
    chip(s, 1.55 + i*1.45, 3.82, 1.35, 0.24, c, PURPLE, RGBColor(0xF0,0xEF,0xFA), 8)
txt(s, "Data Access", 6.05, 3.84, 1.2, 0.22, size=8, bold=True, color=MUTED, font="Calibri")
chip(s, 7.0, 3.82, 2.6, 0.24, "EF Core  ·  AppDbContext", AMBER, RGBColor(0xFB,0xF1,0xE7), 8)
# auth note
box(s, 0.45, 4.10, 9.1, 0.18, RGBColor(0xFB,0xF4,0xE8))
txt(s, "🔐  ทุก endpoint เขียนข้อมูล (master data) ต้องมี JWT + Role = Admin  ·  PBKDF2 hash รหัสผ่าน",
    0.55, 4.10, 9.0, 0.18, size=7.5, color=ORANGED, vmid=True, font="Calibri")

darrow(s, 2.7, 4.40, clr=TEAL)
txt(s, "อ่าน/เขียนข้อมูล", 3.05, 4.44, 2.0, 0.22, size=8, color=TEAL, italic=True, vmid=True, font="Calibri")
darrow(s, 7.1, 4.40, clr=GREY)
txt(s, "serve รูป /assets/*", 7.45, 4.44, 2.2, 0.22, size=8, color=GREY, italic=True, vmid=True, font="Calibri")

# ── Tier 4: Data + Storage ──
box(s, 0.3, 4.80, 4.55, 0.66, RGBColor(0xE9,0xF6,0xF4), TEAL, 1.0, radius=True)
txt(s, "🗄  Database", 0.45, 4.83, 3, 0.24, size=9.5, bold=True, color=TEAL, font="Calibri")
txt(s, "SQLite (AtmInventory.db) ปัจจุบัน  →  MySQL (production)", 0.45, 5.08, 4.3, 0.3, size=8, color=TEXT, font="Calibri")

box(s, 5.15, 4.80, 4.55, 0.66, RGBColor(0xF1,0xF2,0xF4), GREY, 1.0, radius=True)
txt(s, "🖼  Asset Storage (นอก repo)", 5.3, 4.83, 4, 0.24, size=9.5, bold=True, color=GREY, font="Calibri")
txt(s, "D:\\ATMAssets\\parts\\  (529 รูป)  →  serve ที่ /assets/parts/", 5.3, 5.08, 4.3, 0.3, size=8, color=TEXT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ปัญหาของ design เดิม
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "ปัญหา — ทำไมทีมดูตาราง Part แล้วงง", "ตาราง Part เดิม ปนกัน 3 ความหมายในตารางเดียว")

box(s, 0.3, 1.05, 4.6, 3.9, WHITE, GREY_L, 0.75, radius=True)
txt(s, "ตาราง Part เดิม (ปนกัน)", 0.5, 1.15, 4.2, 0.3, size=13, bold=True, color=RED, font="Calibri")
mixed = [
    (BLUE,  "PartNo, PartName, Category, ImagePath", "ตัวอะไหล่ (คืออะไร)"),
    (BLUE,  "MinStock, MaxStock, ReorderPoint", "นโยบายของอะไหล่"),
    (TEAL,  "StockQuantity", "จำนวนในคลัง (เลขเดียว!)"),
    (AMBER, "SerialNo, ExpiryDate, IsUnrepairable", "ชิ้นเฉพาะตัว"),
    (AMBER, "Aging", "ชิ้นเฉพาะตัว (คำนวณได้)"),
    (GREY,  "OrderNumber, TrackingNumber", "ข้อมูลตอนรับเข้า"),
]
yy = 1.55
for clr, fld, mean in mixed:
    box(s, 0.5, yy+0.03, 0.14, 0.42, clr)
    txt(s, fld, 0.75, yy, 2.55, 0.5, size=8.5, bold=True, color=TEXT, vmid=True)
    txt(s, mean, 3.35, yy, 1.5, 0.5, size=8, color=MUTED, vmid=True, italic=True, font="Calibri")
    yy += 0.55

box(s, 5.1, 1.05, 4.6, 3.9, REDBG, RED, 0.75, radius=True)
txt(s, "ผลที่ตามมา", 5.3, 1.15, 4.2, 0.3, size=13, bold=True, color=RED, font="Calibri")
probs = [
    "StockQuantity เป็นเลขเดียว → ของอยู่หลายคลัง\n   (DHL, ราชบูรณะ, ช่างถือ) ตอบไม่ได้ว่าอยู่ไหน",
    "SerialNo อยู่บน Part → 1 รุ่นมีหลายชิ้น\n   หลายซีเรียล แต่ใส่ได้ค่าเดียว",
    "ExpiryDate อยู่บน Part → คนละ lot\n   หมดอายุไม่พร้อมกัน เก็บค่าเดียวผิด",
    "ทีมอ่านแล้วงง: ตกลงแถวนี้คือ \"ชนิดอะไหล่\"\n   หรือ \"ของจริง 1 ชิ้น\" ?",
]
yy = 1.6
for p in probs:
    txt(s, "✕", 5.3, yy, 0.3, 0.7, size=12, bold=True, color=RED, font="Calibri")
    txt(s, p, 5.62, yy, 3.95, 0.7, size=9, color=TEXT, font="Calibri")
    yy += 0.82

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — แนวคิดใหม่ 3 ชั้น (overview)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "แนวคิดใหม่ — แยกเป็น 3 ชั้น + 1 ประวัติ", "แต่ละตารางมีหน้าที่เดียว ชัดเจน")

def layer(slide, x, y, w, h, clr, bgc, name, th, sub, rows):
    box(slide, x, y, w, h, bgc, clr, 1.0, radius=True)
    box(slide, x, y, w, 0.36, clr, radius=True)
    box(slide, x, y+0.18, w, 0.18, clr)
    txt(slide, name, x+0.12, y+0.02, w-0.2, 0.32, size=12, bold=True, color=WHITE, vmid=True, font="Calibri")
    txt(slide, th, x+0.12, y+0.42, w-0.2, 0.26, size=10, bold=True, color=clr, font="Calibri")
    txt(slide, sub, x+0.12, y+0.66, w-0.2, 0.4, size=8.5, color=MUTED, italic=True, font="Calibri")
    yy = y+1.04
    for r in rows:
        txt(slide, "• "+r, x+0.14, yy, w-0.25, 0.24, size=8.5, color=TEXT, font="Calibri")
        yy += 0.235

layer(s, 0.3, 1.05, 2.32, 3.5, BLUE, RGBColor(0xEE,0xF3,0xFE), "Part", "ตัวอะไหล่",
      "1 แถว = 1 รหัสอะไหล่", ["คืออะไร ชื่อ หมวด", "รูป ต้นทุน หน่วย", "min/max/reorder", "ไม่มีจำนวน", "ไม่มีซีเรียล"])
layer(s, 2.78, 1.05, 2.32, 3.5, TEAL, RGBColor(0xE9,0xF6,0xF4), "PartStock", "จำนวนในคลัง",
      "1 แถว = อะไหล่ x คลัง", ["GoodQty (ดี)", "DefectiveQty (เสีย)", "ผูก Part + Location", "UNIQUE(Part,Loc)", "ยอดคงเหลือจริง"])
layer(s, 5.26, 1.05, 2.32, 3.5, AMBER, RGBColor(0xFB,0xF1,0xE7), "PartUnit", "ชิ้นเฉพาะตัว",
      "1 แถว = 1 ชิ้นจริง (option)", ["SerialNo (unique)", "ExpiryDate", "IsUnrepairable", "Condition", "เฉพาะของ track serial"])
layer(s, 7.74, 1.05, 1.96, 3.5, GREY, RGBColor(0xF1,0xF2,0xF4), "StockMovement", "ประวัติ",
      "1 แถว = 1 การขยับ", ["GR/Issue/Return", "Transfer/Disposal", "From/To Location", "ใครทำ เมื่อไหร่", "ใช้ทำ Audit"])

# arrows between layers
def arr(slide, x1, x2, y, clr):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = clr; c.line.width = Pt(1.4); c.shadow.inherit = False
txt(s, "1 : N", 2.5, 2.6, 0.6, 0.22, size=8.5, bold=True, color=MUTED, font="Calibri", align=PP_ALIGN.CENTER)
txt(s, "1 : N", 4.98, 2.6, 0.6, 0.22, size=8.5, bold=True, color=MUTED, font="Calibri", align=PP_ALIGN.CENTER)

txt(s, "อยากรู้ \"เหลือกี่ชิ้นทั้งหมด\" ?  →  SUM(PartStock.GoodQty) ของ Part นั้น  (ไม่เก็บเป็นเลขนิ่งใน Part อีกต่อไป)",
    0.3, 4.78, 9.4, 0.4, size=9.5, color=TEAL, bold=True, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ERD diagram (boxes + crow's-foot relationships)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "ERD — แผนภาพความสัมพันธ์ทั้งหมด", "กล่อง = ตาราง  ·  เส้น = ความสัมพันธ์ (1 : N)")

def entity(slide, x, y, w, name, sub, clr, bgc, fields):
    """Draw an ERD entity box with a colored title bar + field rows."""
    rows_h = len(fields)*0.205
    h = 0.5 + rows_h
    box(slide, x, y, w, h, bgc, clr, 1.25, radius=True)
    box(slide, x, y, w, 0.32, clr)
    txt(slide, name, x+0.1, y+0.01, w-0.15, 0.3, size=10.5, bold=True, color=WHITE, vmid=True, font="Calibri")
    txt(slide, sub, x+0.1, y+0.33, w-0.15, 0.18, size=7, color=clr, bold=True, font="Calibri")
    yy = y + 0.52
    for ic, fn in fields:
        if ic:
            kc = {"PK":ORANGED,"UK":PURPLE,"FK":GREEN}.get(ic, MUTED)
            txt(slide, ic, x+0.08, yy, 0.32, 0.2, size=7, bold=True, color=kc, vmid=True, font="Calibri")
        txt(slide, fn, x+(0.42 if ic else 0.12), yy, w-0.5, 0.2, size=8, color=TEXT, vmid=True)
        yy += 0.205
    return (x, y, w, h)

def crow(slide, x1, y1, x2, y2, clr=MUTED):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = clr; c.line.width = Pt(1.5); c.shadow.inherit = False

# Master entities (top)
entity(s, 0.3, 1.0, 1.95, "Category", "MASTER", BLUE, RGBColor(0xEE,0xF3,0xFE),
       [("PK","Id"),("UK","Name"),("","Description")])
entity(s, 7.55, 1.0, 2.1, "Location", "MASTER", TEAL, RGBColor(0xE9,0xF6,0xF4),
       [("PK","Id"),("UK","Code"),("","Name · Type")])

# Part (center-left)
px,py,pw,ph = entity(s, 0.3, 2.55, 2.25, "Part", "ตัวอะไหล่", BLUE, RGBColor(0xEE,0xF3,0xFE),
       [("PK","Id"),("UK","PartNo"),("FK","CategoryId"),("","PartName"),
        ("","Unit · MainUnit"),("","CostPerUnit"),("","ImagePath"),("","min/max/reorder")])

# PartStock (center)
entity(s, 3.15, 2.55, 2.05, "PartStock", "ในคลัง", TEAL, RGBColor(0xE9,0xF6,0xF4),
       [("PK","Id"),("FK","PartId"),("FK","LocationId"),("","GoodQty"),
        ("","DefectiveQty"),("","UQ(Part,Loc)")])

# PartUnit (right)
entity(s, 5.6, 2.55, 2.05, "PartUnit", "ชิ้นเฉพาะ", AMBER, RGBColor(0xFB,0xF1,0xE7),
       [("PK","Id"),("FK","PartId"),("FK","LocationId"),("UK","SerialNo"),
        ("","ExpiryDate"),("","Condition")])

# StockMovement (bottom-right, ledger)
entity(s, 7.95, 2.55, 1.75, "StockMovement", "LEDGER", GREY, RGBColor(0xF1,0xF2,0xF4),
       [("PK","Id"),("FK","PartId"),("","From/To Loc"),("","Type · Qty"),
        ("","User · Time")])

# Relationships — routed via top/bottom buses so no line crosses a box
# 1) Category 1:N Part  (vertical)
crow(s, 1.27, 2.115, 1.27, 2.55, BLUE)
txt(s, "1:N", 1.32, 2.24, 0.6, 0.2, size=8, bold=True, color=BLUE, font="Calibri")
# 2) Part 1:N PartStock  (short horizontal, mid-height)
crow(s, 2.55, 3.4, 3.15, 3.4, TEAL)
txt(s, "1:N", 2.62, 3.2, 0.6, 0.2, size=8, bold=True, color=TEAL, font="Calibri")
# 3) Location 1:N PartStock  (top bus at y=2.40, above the boxes)
crow(s, 7.70, 2.115, 7.70, 2.40, TEAL)
crow(s, 7.70, 2.40, 4.175, 2.40, TEAL)
crow(s, 4.175, 2.40, 4.175, 2.55, TEAL)
txt(s, "1:N", 5.85, 2.18, 0.6, 0.2, size=8, bold=True, color=TEAL, font="Calibri")
# 4) + 5) Part 1:N PartUnit  &  Part 1:N StockMovement  (bottom bus at y=4.90)
crow(s, 1.425, 4.69, 1.425, 4.90, GREY)
crow(s, 1.425, 4.90, 8.825, 4.90, GREY)
crow(s, 6.625, 4.90, 6.625, 4.28, AMBER)         # riser to PartUnit
crow(s, 8.825, 4.90, 8.825, 4.075, GREY)         # riser to StockMovement
txt(s, "1:N (option)", 6.05, 4.93, 1.5, 0.2, size=8, bold=True, color=AMBER, font="Calibri")
txt(s, "1:N (ทุกการขยับ)", 2.4, 4.70, 2.2, 0.2, size=8, bold=True, color=GREY, font="Calibri")

legend(s, 5.30)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Part (Master) after cleanup
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Part — ตัวอะไหล่ (Master / Catalog)", "เหลือเฉพาะ \"นิยามของอะไหล่\" — สะอาด อ่านเข้าใจทันที")
part = [
    ("Id", "int", False, "PK", "Primary key — auto increment"),
    ("PartNo", "varchar(50)", False, "UK", "รหัสอะไหล่ — business key (ห้ามซ้ำ)"),
    ("PartName", "varchar(255)", False, "", "ชื่อ/คำอธิบายอะไหล่"),
    ("CategoryId", "int", True, "FK", "→ Category (Sub Unit)"),
    ("Unit", "varchar(20)", False, "", "หน่วยนับ default \"pcs\""),
    ("MainUnit", "varchar(100)", True, "", "กลุ่มหลัก เช่น Cabinet"),
    ("CostPerUnit", "decimal(18,2)", True, "", "ต้นทุนต่อหน่วย"),
    ("CatalogueRef", "varchar(100)", True, "", "อ้างอิงเลขใน catalogue"),
    ("ImagePath", "varchar(255)", True, "", "พาธรูป /assets/parts/..."),
    ("Remark", "varchar(500)", True, "", "หมายเหตุอิสระ"),
    ("MinStock", "int", False, "", "สต็อกขั้นต่ำ (นโยบาย)"),
    ("MaxStock", "int", False, "", "สต็อกขั้นสูง (นโยบาย)"),
    ("ReorderPoint", "int", False, "", "จุดสั่งซื้อใหม่ (นโยบาย)"),
    ("IsActive", "tinyint(1)", False, "", "สถานะ — false = soft delete"),
    ("CreatedAt / UpdatedAt", "datetime", False, "", "เวลาสร้าง / แก้ไขล่าสุด"),
]
field_table(s, 0.3, 1.05, part, col_w=(1.85,1.45,0.45,0.45,5.5), row_h=0.262, fsize=8.5)
legend(s, 5.34)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PartStock (Balance)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "PartStock — จำนวนในคลัง (Balance)", "ยอดคงเหลือจริง แยกตามคลัง + สภาพดี/เสีย")
ps = [
    ("Id", "int", False, "PK", "Primary key"),
    ("PartId", "int", False, "FK", "→ Part — อะไหล่ตัวไหน"),
    ("LocationId", "int", False, "FK", "→ Location — คลังไหน"),
    ("GoodQty", "int", False, "", "จำนวนสภาพดี (พร้อมใช้)"),
    ("DefectiveQty", "int", False, "", "จำนวนสภาพเสีย (รอซ่อม/ทิ้ง)"),
    ("UpdatedAt", "datetime", False, "", "เวลาปรับยอดล่าสุด"),
]
field_table(s, 0.3, 1.2, ps, col_w=(1.6,1.5,0.5,0.5,5.05), row_h=0.36, fsize=9.5)

box(s, 0.3, 3.7, 9.4, 1.25, RGBColor(0xE9,0xF6,0xF4), TEAL, 0.75, radius=True)
txt(s, "กฎสำคัญ", 0.5, 3.78, 3, 0.3, size=11, bold=True, color=TEAL, font="Calibri")
txt(s, "UNIQUE (PartId, LocationId)  — อะไหล่ 1 ตัว ต่อ 1 คลัง มีได้แถวเดียว (กันยอดซ้ำ/ตีกัน)",
    0.5, 4.08, 9.0, 0.3, size=10, color=TEXT, font="Calibri")
txt(s, "ยอดรวมทั้งระบบของอะไหล่ตัวหนึ่ง  =  SUM(GoodQty) ทุก Location  — แทน Part.StockQuantity เดิม",
    0.5, 4.42, 9.0, 0.3, size=10, color=TEXT, font="Calibri")
txt(s, "ทุกครั้งที่ยอดเปลี่ยน ต้องผ่าน StockService → เขียน StockMovement 1 แถวเสมอ (กันยอดเพี้ยน)",
    0.5, 4.72, 9.0, 0.3, size=10, color=TEXT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PartUnit (Serial, optional)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "PartUnit — ชิ้นเฉพาะตัว (Serial · ทางเลือก)", "ใช้เฉพาะอะไหล่ที่ต้อง track ทีละชิ้น เช่น มีซีเรียล / มีวันหมดอายุ")
pu = [
    ("Id", "int", False, "PK", "Primary key"),
    ("PartId", "int", False, "FK", "→ Part — เป็นอะไหล่รุ่นไหน"),
    ("LocationId", "int", True, "FK", "→ Location — ชิ้นนี้อยู่คลังไหน"),
    ("SerialNo", "varchar(100)", False, "UK", "หมายเลขซีเรียลของชิ้นนี้ (ห้ามซ้ำ)"),
    ("Condition", "varchar(20)", False, "", "Good / Defective"),
    ("ExpiryDate", "datetime", True, "", "วันหมดอายุของชิ้นนี้ (ตาม lot)"),
    ("IsUnrepairable", "tinyint(1)", False, "", "ซ่อมไม่ได้ → เข้า Disposal"),
    ("ReceivedAt", "datetime", False, "", "วันรับเข้า — ใช้คำนวณ Aging"),
    ("Status", "varchar(20)", False, "", "InStock / Issued / Disposed"),
]
field_table(s, 0.3, 1.05, pu, col_w=(1.65,1.5,0.5,0.5,5.0), row_h=0.295, fsize=9)

box(s, 0.3, 4.2, 9.4, 0.78, RGBColor(0xFB,0xF1,0xE7), AMBER, 0.75, radius=True)
txt(s, "เมื่อไหร่ถึงใช้ตารางนี้", 0.5, 4.28, 4, 0.3, size=11, bold=True, color=AMBER, font="Calibri")
txt(s, "อะไหล่ทั่วไป (น็อต สายไฟ) → ไม่ต้องมี PartUnit ใช้แค่ PartStock นับจำนวนพอ   |   "
       "อะไหล่ที่ต้องรู้รายชิ้น (mainboard มีซีเรียล, ของมีวันหมดอายุ) → สร้าง PartUnit 1 แถวต่อชิ้น",
    0.5, 4.56, 9.0, 0.4, size=9, color=TEXT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — StockMovement (Ledger)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "StockMovement — ประวัติทุกการเคลื่อนไหว (Ledger)", "เขียนทุกครั้งที่ของขยับ — ห้ามแก้ ห้ามลบ (append-only)")
sm = [
    ("Id", "int", False, "PK", "Primary key"),
    ("MovementType", "varchar(20)", False, "", "GR / Issue / Return / Transfer / Disposal / Adjust"),
    ("PartId", "int", False, "FK", "→ Part"),
    ("FromLocationId", "int", True, "FK", "→ Location ต้นทาง (ออก)"),
    ("ToLocationId", "int", True, "FK", "→ Location ปลายทาง (เข้า)"),
    ("Qty", "int", False, "", "จำนวนที่ขยับ"),
    ("Condition", "varchar(20)", False, "", "Good / Defective"),
    ("RefType", "varchar(20)", True, "", "Ticket / GoodsReceipt / Transfer ..."),
    ("RefId", "varchar(50)", True, "", "เลขเอกสารอ้างอิง"),
    ("Cost", "decimal(18,2)", True, "", "ต้นทุน ณ ตอนขยับ"),
    ("UserName", "varchar(100)", False, "", "ใครเป็นคนทำ"),
    ("Timestamp", "datetime", False, "", "เวลาที่ทำรายการ"),
]
field_table(s, 0.3, 1.05, sm, col_w=(1.7,1.5,0.45,0.45,5.6), row_h=0.275, fsize=8.5)
txt(s, "ทำไมสำคัญ:  ใช้ทำ Audit Report, คำนวณ Aging, Lifecycle ของอะไหล่ — เป็นหลักฐานว่ายอดในคลังมาจากไหน",
    0.3, 5.0, 9.4, 0.3, size=9.5, color=GREY, bold=True, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Before / After mapping
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Before → After — ฟิลด์ไหน ย้ายไปไหน", "สรุปการเปลี่ยนแปลงให้ทีมเห็นภาพรวม")

# table header
cols_x = [0.3, 3.0, 6.5]
cw     = [2.7, 3.5, 3.2]
box(s, 0.3, 1.05, 9.4, 0.34, NAVY)
for i,(cx,cap) in enumerate(zip(cols_x, ["ฟิลด์เดิมใน Part","ย้ายไปที่","เหตุผล"])):
    txt(s, cap, cx+0.1, 1.07, cw[i]-0.15, 0.3, size=10, bold=True, color=WHITE, vmid=True, font="Calibri")

rows = [
    ("StockQuantity", "ลบทิ้ง → SUM(PartStock)", "ของอยู่หลายคลัง เลขเดียวไม่พอ", RED),
    ("SerialNo", "PartUnit.SerialNo", "1 รุ่นมีหลายซีเรียล", AMBER),
    ("ExpiryDate", "PartUnit.ExpiryDate", "คนละ lot หมดอายุไม่พร้อมกัน", AMBER),
    ("IsUnrepairable", "PartUnit.IsUnrepairable", "เป็นสถานะของชิ้น ไม่ใช่ของรุ่น", AMBER),
    ("Aging", "คำนวณจาก PartUnit.ReceivedAt", "ค่าที่คำนวณได้ ไม่ควรเก็บนิ่ง", GREY),
    ("OrderNumber", "GoodsReceipt", "ข้อมูลตอนรับเข้า (transaction)", GREY),
    ("TrackingNumber", "GoodsReceipt / Transfer", "ข้อมูลการขนส่ง (transaction)", GREY),
]
ry = 1.39
for fld, to, why, clr in rows:
    box(s, 0.3, ry, 9.4, 0.43, WHITE, GREY_L, 0.5)
    box(s, 0.3, ry, 0.08, 0.43, clr)
    txt(s, fld, 0.48, ry, cw[0]-0.2, 0.43, size=10, bold=True, color=TEXT, vmid=True)
    txt(s, "→  "+to, cols_x[1]+0.05, ry, cw[1]-0.1, 0.43, size=9.5, color=clr, bold=True, vmid=True, font="Calibri")
    txt(s, why, cols_x[2]+0.05, ry, cw[2]-0.1, 0.43, size=9, color=MUTED, vmid=True, font="Calibri")
    ry += 0.475

txt(s, "เหลือใน Part เฉพาะ \"นิยามของอะไหล่\" → ทีมเปิดตารางปุ๊บ เข้าใจทันทีว่าคือ catalog",
    0.3, ry+0.03, 9.4, 0.3, size=10, color=GREEN, bold=True, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Relationships & migration plan
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, NAVY)
box(s, 0, 0, 0.55, 5.625, ORANGE)
txt(s, "ความสัมพันธ์ & แผนย้าย", 0.9, 0.45, 8.7, 0.65, size=24, bold=True, color=WHITE, font="Calibri")

rels = [
    ("Category  1 : N  Part", "Part.CategoryId → Category.Id", "SetNull"),
    ("Part  1 : N  PartStock", "PartStock.PartId → Part.Id", "Cascade"),
    ("Location  1 : N  PartStock", "PartStock.LocationId → Location.Id", "Restrict"),
    ("Part  1 : N  PartUnit", "PartUnit.PartId → Part.Id  (optional)", "Cascade"),
    ("Part  1 : N  StockMovement", "StockMovement.PartId → Part.Id", "Restrict"),
]
for j, (rel, via, beh) in enumerate(rels):
    ry = 1.2 + j*0.5
    box(s, 0.9, ry, 8.7, 0.42, RGBColor(0x12,0x28,0x40), RGBColor(0x2A,0x4A,0x70), 0.75)
    txt(s, rel, 1.05, ry, 3.5, 0.42, size=11, bold=True, color=ORANGE, vmid=True, font="Calibri")
    txt(s, via, 4.5, ry, 3.6, 0.42, size=9.5, color=WHITE, vmid=True)
    txt(s, beh, 8.15, ry, 1.4, 0.42, size=9, color=RGBColor(0xB0,0xC8,0xE8), vmid=True, italic=True, font="Calibri")

txt(s, "แผนย้ายแบบไม่เสียข้อมูล (3 ขั้น)", 0.9, 3.85, 6, 0.3, size=13, bold=True, color=ORANGE, font="Calibri")
steps = [
    "1.  สร้าง PartStock จากข้อมูลเดิม: 1 แถวต่อ (Part, คลังหลัก) ใส่ GoodQty = StockQuantity เดิม",
    "2.  เปลี่ยนทุกที่ที่อ่าน Part.StockQuantity ให้ไปอ่าน SUM(PartStock) แทน แล้วค่อยลบคอลัมน์",
    "3.  PartUnit ทำทีหลังได้ — เพิ่มเฉพาะตอนเริ่ม track ซีเรียลจริงจัง (schema เผื่อไว้แล้ว)",
]
for j, st in enumerate(steps):
    txt(s, st, 0.95, 4.2+j*0.37, 8.7, 0.34, size=10, color=WHITE, font="Calibri")

out = "D:/ATMApi/ATM-Inventory-System/ATM_Inventory_ERD_Redesign.pptx"
prs.save(out)
print(f"Saved: {out}")
