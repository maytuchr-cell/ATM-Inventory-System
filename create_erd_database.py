"""
ATM Inventory System — Phase 1 ERD Database (detailed field reference + sizes)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR

NAVY   = RGBColor(0x1C, 0x35, 0x57)
BLUE   = RGBColor(0x0D, 0x7E, 0xD8)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
ORANGED= RGBColor(0xD8, 0x5A, 0x30)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
TEXT   = RGBColor(0x1E, 0x2D, 0x40)
TEAL   = RGBColor(0x02, 0x8C, 0x9A)
GREEN  = RGBColor(0x05, 0x96, 0x69)
RED    = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x53, 0x4A, 0xB7)
GREY_L = RGBColor(0xD1, 0xD5, 0xDB)
ROW_ALT= RGBColor(0xF3, 0xF6, 0xFB)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])
def bg(slide, c):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def box(slide, x, y, w, h, fill, line=None, lw=0.75, radius=False):
    sh = slide.shapes.add_shape(5 if radius else 1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(slide, text, x, y, w, h, size=12, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, vmid=False, font="Consolas"):
    if color is None: color = TEXT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.03); tf.margin_right=Inches(0.03); tf.margin_top=Inches(0.01); tf.margin_bottom=Inches(0.01)
    if vmid: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def header(slide, title, subtitle=None):
    box(slide, 0, 0, 10, 0.78, NAVY)
    box(slide, 0, 0.78, 10, 0.04, ORANGE)
    txt(slide, title, 0.4, 0.08, 9.2, 0.55, size=20, bold=True, color=WHITE, vmid=True, font="Calibri")
    if subtitle:
        txt(slide, subtitle, 0.42, 0.88, 9, 0.26, size=10, color=MUTED, italic=True, font="Calibri")

# Field table: columns Field | Data Type (size) | Null | Key | Description
def field_table(slide, x, y, fields, col_w=(1.7, 1.55, 0.5, 0.5, 4.45), row_h=0.265, fsize=8.5, header_clr=NAVY):
    cols = ["Field", "Data Type", "Null", "Key", "Description"]
    tw = sum(col_w)
    box(slide, x, y, tw, 0.3, header_clr)
    cx = x
    for i, c in enumerate(cols):
        al = PP_ALIGN.CENTER if c in ("Null","Key") else PP_ALIGN.LEFT
        txt(slide, c, cx+0.05, y+0.02, col_w[i]-0.08, 0.26, size=8.5, bold=True, color=WHITE, vmid=True, font="Calibri", align=al)
        cx += col_w[i]
    box(slide, x, y+0.3, tw, len(fields)*row_h, WHITE, GREY_L, 0.75)
    for ri, (fname, ftype, nullable, key, desc) in enumerate(fields):
        ry = y + 0.3 + ri*row_h
        if ri % 2 == 1: box(slide, x+0.02, ry, tw-0.04, row_h, ROW_ALT)
        cx = x
        txt(slide, fname, cx+0.06, ry, col_w[0]-0.1, row_h, size=fsize, bold=True, color=TEXT, vmid=True)
        cx += col_w[0]
        txt(slide, ftype, cx+0.06, ry, col_w[1]-0.1, row_h, size=fsize-0.5, color=BLUE, vmid=True)
        cx += col_w[1]
        nl = "NULL" if nullable else "NOT"
        nc = MUTED if nullable else RGBColor(0x9a,0x9a,0x9a)
        txt(slide, nl, cx+0.02, ry, col_w[2]-0.04, row_h, size=7, color=nc, vmid=True, align=PP_ALIGN.CENTER, font="Calibri")
        cx += col_w[2]
        kc = {"PK":ORANGED, "UK":PURPLE, "FK":GREEN}.get(key, MUTED)
        if key: txt(slide, key, cx+0.03, ry, col_w[3]-0.06, row_h, size=fsize-0.5, bold=True, color=kc, vmid=True, align=PP_ALIGN.CENTER, font="Calibri")
        cx += col_w[3]
        txt(slide, desc, cx+0.06, ry, col_w[4]-0.1, row_h, size=fsize, color=TEXT, vmid=True, font="Calibri")

def legend(slide, y=5.3):
    items = [("PK", ORANGED, "Primary"), ("UK", PURPLE, "Unique"), ("FK", GREEN, "Foreign")]
    x = 0.4
    for tag, clr, desc in items:
        txt(slide, tag, x, y-0.02, 0.4, 0.25, size=9, bold=True, color=clr, vmid=True, font="Calibri"); x += 0.34
        txt(slide, desc, x, y-0.02, 1.2, 0.25, size=9, color=MUTED, vmid=True, font="Calibri"); x += len(desc)*0.06 + 0.25
    txt(slide, "Null = nullable column", x, y-0.02, 2.2, 0.25, size=9, color=MUTED, vmid=True, font="Calibri"); x += 1.6
    box(slide, x, y, 0.4, 0.2, ORANGED); txt(slide, "NEW", x, y, 0.4, 0.2, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, vmid=True, font="Calibri"); x += 0.5
    txt(slide, "added from GRG catalog", x, y-0.02, 2.5, 0.25, size=9, color=MUTED, vmid=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, NAVY)
box(s, 0, 0, 0.55, 5.625, ORANGE)
txt(s, "ERD — Database Design", 0.95, 1.05, 8.5, 1.0, size=38, bold=True, color=WHITE, font="Calibri")
txt(s, "Phase 1 Master Data — Detailed Field Reference", 0.97, 2.1, 8.5, 0.5, size=18, color=RGBColor(0xB0,0xC8,0xE8), font="Calibri")
box(s, 0.97, 2.85, 7.3, 0.045, ORANGE)
txt(s, "Part Master · Categories · Part Configuration", 0.97, 3.0, 8, 0.4, size=14, color=RGBColor(0x90,0xAE,0xD8), font="Calibri")
txt(s, "ATM Inventory System  |  DataOne Asia  |  June 2026", 0.97, 3.72, 8, 0.36, size=12, color=MUTED, italic=True, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ERD Overview
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "ERD Overview — ความสัมพันธ์ระหว่างตาราง", "5 ตารางหลักของ Master Data Phase 1")

def mini(slide, x, y, w, name, clr, rows):
    box(slide, x, y, w, 0.34, clr)
    txt(slide, name, x+0.1, y+0.03, w-0.15, 0.28, size=11, bold=True, color=WHITE, vmid=True, font="Calibri")
    box(slide, x, y+0.34, w, len(rows)*0.26, WHITE, clr, 1.0)
    for j, (ic, fn) in enumerate(rows):
        ry = y+0.34+j*0.26
        if ic:
            kc = {"PK":ORANGED,"UK":PURPLE,"FK":GREEN,"→":PURPLE}.get(ic, MUTED)
            txt(slide, ic, x+0.07, ry, 0.4, 0.26, size=8, bold=True, color=kc, vmid=True, font="Calibri")
        txt(slide, fn, x+(0.5 if ic else 0.12), ry, w-0.55, 0.26, size=9, color=TEXT, vmid=True)

mini(s, 0.3, 1.45, 1.85, "Category", GREEN, [("PK","id"),("","name (Sub Unit)"),("","isActive")])
mini(s, 0.3, 3.2, 1.85, "Vendor*", MUTED, [("PK","id"),("","code · type")])
mini(s, 3.1, 1.25, 2.3, "Part", PURPLE, [("PK","id"),("UK","partNo"),("FK","categoryId"),("","mainUnit / remark"),("","imagePath"),("","stock / min / max")])
mini(s, 6.55, 1.2, 2.0, "AtmModel (ATM Group)", BLUE, [("PK","id"),("","modelCode"),("→","partNo")])
mini(s, 6.55, 3.0, 2.0, "EquivalentGroup", TEAL, [("PK","id"),("","name"),("→","partNo")])

def conn(slide, x1,y1,x2,y2,clr):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=clr; c.line.width=Pt(1.6); c.shadow.inherit=False
conn(s, 2.15,1.9,3.1,2.0, GREEN); txt(s,"1 : N",2.3,1.62,0.8,0.22,size=9,bold=True,color=GREEN,font="Calibri"); txt(s,"categoryId",2.25,2.04,1.0,0.22,size=8,color=MUTED,italic=True,font="Calibri")
conn(s, 5.4,1.7,6.55,1.7, PURPLE); txt(s,"N : M  partNo",5.45,1.48,1.3,0.22,size=8.5,bold=True,color=PURPLE,font="Calibri")
conn(s, 5.4,2.2,6.55,3.4, PURPLE); txt(s,"N : M  partNo",5.4,2.9,1.3,0.22,size=8.5,bold=True,color=PURPLE,font="Calibri")

txt(s, "* Vendor เป็น master data แต่ใช้ในโมดูล Goods Receipt (Phase 2) — ไม่ผูกกับ Part โดยตรง",
    0.3, 4.95, 9.4, 0.3, size=9, color=MUTED, italic=True, font="Calibri")
legend(s, 5.32)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Part (identity + catalog)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Part — Identity & Catalog Fields", "ตาราง Parts (1/2) — ชนิดข้อมูลและขนาด")

part1 = [
    ("Id", "int", False, "PK", "Primary key — auto increment"),
    ("OrderNumber", "varchar(50)", False, "", "เลขที่สั่งซื้อ/อ้างอิงเอกสาร"),
    ("PartNo", "varchar(50)", False, "UK", "รหัสอะไหล่ — business key (Part Number)"),
    ("PartName", "varchar(255)", False, "", "ชื่อ/คำอธิบายอะไหล่ (Part Description)"),
    ("SerialNo", "varchar(100)", True, "", "หมายเลข Serial (ถ้า track serial)"),
    ("CatalogueRef", "varchar(100)", True, "", "อ้างอิงหมายเลขใน catalogue"),
    ("CategoryId", "int", True, "FK", "→ Category — กลุ่มย่อย (Sub Unit)"),
    ("MainUnit", "varchar(100)", True, "", "กลุ่มหลัก เช่น Cabinet (Main Unit) — NEW"),
    ("Remark", "varchar(500)", True, "", "หมายเหตุอิสระจาก catalog — NEW"),
    ("ImagePath", "varchar(255)", True, "", "พาธรูป /uploads/parts/... (Picture) — NEW"),
    ("Unit", "varchar(20)", False, "", "หน่วยนับ default \"pcs\""),
]
field_table(s, 0.3, 1.05, part1, row_h=0.30, fsize=8.5)
for ny in (1.05+0.3+7*0.30, 1.05+0.3+8*0.30, 1.05+0.3+9*0.30):
    box(s, 8.95, ny+0.04, 0.34, 0.18, ORANGED); txt(s, "NEW", 8.95, ny+0.04, 0.34, 0.18, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER, vmid=True, font="Calibri")
legend(s, 5.32)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Part (stock + lifecycle)
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Part — Stock, Cost & Lifecycle Fields", "ตาราง Parts (2/2) — สต็อก ต้นทุน และสถานะ")

part2 = [
    ("StockQuantity", "int", False, "", "จำนวนคงเหลือรวมทุกคลัง (denormalized)"),
    ("MinStock", "int", False, "", "สต็อกขั้นต่ำ — Below Min (default 1)"),
    ("MaxStock", "int", False, "", "สต็อกขั้นสูง — Over Max (default 100)"),
    ("ReorderPoint", "int", False, "", "จุดสั่งซื้อใหม่ — Reorder (default 3)"),
    ("CostPerUnit", "decimal(18,2)", True, "", "ต้นทุนต่อหน่วย"),
    ("Aging", "int", True, "", "อายุสต็อก (จำนวนวัน)"),
    ("TrackingNumber", "varchar(100)", True, "", "เลขติดตามการขนส่ง"),
    ("ExpiryDate", "datetime", True, "", "วันหมดอายุ (โมดูล Disposal)"),
    ("IsUnrepairable", "tinyint(1)", False, "", "ทำเครื่องหมายซ่อมไม่ได้ (Disposal)"),
    ("IsActive", "tinyint(1)", False, "", "สถานะใช้งาน — false = soft delete"),
    ("CreatedAt", "datetime", False, "", "เวลาสร้างเรคคอร์ด (UTC)"),
    ("UpdatedAt", "datetime", False, "", "เวลาแก้ไขล่าสุด (UTC)"),
]
field_table(s, 0.3, 1.05, part2, row_h=0.275, fsize=8.5)
legend(s, 5.32)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Category
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Category — หมวดหมู่อะไหล่ (Sub Unit)", "ตาราง Categories")

cat = [
    ("Id", "int", False, "PK", "Primary key — auto increment"),
    ("Name", "varchar(100)", False, "UK", "ชื่อหมวดหมู่ = Sub Unit (ห้ามซ้ำ)"),
    ("Description", "varchar(500)", True, "", "รายละเอียดหมวดหมู่"),
    ("IsActive", "tinyint(1)", False, "", "สถานะ — ลบไม่ได้ถ้ามี Part active"),
    ("Parts", "navigation", False, "", "ความสัมพันธ์ 1:N → Part"),
    ("partCount", "int (computed)", False, "", "จำนวน Part active (คำนวณที่ API)"),
]
field_table(s, 0.3, 1.2, cat, col_w=(1.55,1.65,0.5,0.5,5.0), row_h=0.42, fsize=10)

box(s, 0.3, 4.0, 9.4, 1.05, RGBColor(0xEA,0xF6,0xF0), GREEN, 0.75, radius=True)
txt(s, "ความสัมพันธ์", 0.5, 4.08, 3, 0.3, size=11, bold=True, color=GREEN, font="Calibri")
txt(s, "Part.CategoryId → Category.Id   (OnDelete: SetNull)  — เมื่อลบหมวด อะไหล่ไม่ถูกลบ แต่ categoryId กลายเป็น null",
    0.5, 4.4, 9.0, 0.6, size=10.5, color=TEXT, font="Calibri")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ATM Group
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Part Configuration — ATM Group", "ตาราง AtmModel + AtmModelPart")

txt(s, "AtmModel", 0.35, 0.98, 4, 0.3, size=13, bold=True, color=BLUE, font="Calibri")
atm = [
    ("Id", "int", False, "PK", "Primary key"),
    ("ModelCode", "varchar(50)", False, "", "รหัสกลุ่ม/รุ่น (Group Code) เช่น GRG-H22N"),
    ("ModelName", "varchar(150)", False, "", "ชื่อกลุ่ม/รุ่น (Group Name)"),
    ("Manufacturer", "varchar(100)", True, "", "ผู้ผลิต เช่น GRG, NCR"),
    ("Description", "varchar(500)", True, "", "รายละเอียด"),
    ("IsActive", "tinyint(1)", False, "", "สถานะใช้งาน"),
    ("CompatibleParts", "navigation", False, "", "1:N → AtmModelPart"),
]
field_table(s, 0.3, 1.28, atm, col_w=(1.6,1.5,0.5,0.5,5.0), row_h=0.245, fsize=8.5)

txt(s, "AtmModelPart  (ตารางเชื่อม M:N)", 0.35, 3.32, 5, 0.3, size=13, bold=True, color=BLUE, font="Calibri")
atmp = [
    ("Id", "int", False, "PK", "Primary key"),
    ("AtmModelId", "int", False, "FK", "→ AtmModel (OnDelete: Cascade)"),
    ("PartNo", "varchar(50)", False, "", "อ้างอิง Part.PartNo — UK(AtmModelId, PartNo)"),
]
field_table(s, 0.3, 3.62, atmp, col_w=(1.6,1.5,0.5,0.5,5.0), row_h=0.245, fsize=8.5)
legend(s, 5.32)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Equivalent Group
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, LIGHT)
header(s, "Part Configuration — Equivalent Group", "ตาราง EquivalentGroup + Member")

txt(s, "EquivalentGroup", 0.35, 0.98, 4, 0.3, size=13, bold=True, color=TEAL, font="Calibri")
eg = [
    ("Id", "int", False, "PK", "Primary key"),
    ("Name", "varchar(150)", False, "", "ชื่อกลุ่มอะไหล่ทดแทน"),
    ("Description", "varchar(500)", True, "", "รายละเอียด"),
    ("CreatedAt", "datetime", False, "", "เวลาสร้างกลุ่ม"),
    ("Members", "navigation", False, "", "1:N → EquivalentGroupMember"),
]
field_table(s, 0.3, 1.28, eg, col_w=(1.65,1.5,0.5,0.5,4.95), row_h=0.3, fsize=9)

txt(s, "EquivalentGroupMember  (สมาชิกกลุ่ม)", 0.35, 3.28, 5, 0.3, size=13, bold=True, color=TEAL, font="Calibri")
egm = [
    ("Id", "int", False, "PK", "Primary key"),
    ("GroupId", "int", False, "FK", "→ EquivalentGroup (OnDelete: Cascade)"),
    ("PartNo", "varchar(50)", False, "", "อ้างอิง Part.PartNo — UK(GroupId, PartNo)"),
]
field_table(s, 0.3, 3.58, egm, col_w=(1.65,1.5,0.5,0.5,4.95), row_h=0.3, fsize=9)
legend(s, 5.35)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Relationships & Constraints
# ════════════════════════════════════════════════════════════════════════
s = blank(); bg(s, NAVY)
box(s, 0, 0, 0.55, 5.625, ORANGE)
txt(s, "Relationships & Constraints", 0.9, 0.5, 8.7, 0.65, size=26, bold=True, color=WHITE, font="Calibri")

rels = [
    ("Category  1 : N  Part", "Part.CategoryId → Category.Id", "OnDelete SetNull"),
    ("AtmModel  M : N  Part", "ผ่าน AtmModelPart.PartNo", "OnDelete Cascade"),
    ("EquivalentGroup  M : N  Part", "ผ่าน EquivalentGroupMember.PartNo", "OnDelete Cascade"),
]
for j, (rel, via, beh) in enumerate(rels):
    ry = 1.35 + j*0.62
    box(s, 0.9, ry, 8.7, 0.5, RGBColor(0x12,0x28,0x40), RGBColor(0x2A,0x4A,0x70), 0.75)
    txt(s, rel, 1.05, ry, 3.6, 0.5, size=12, bold=True, color=ORANGE, vmid=True, font="Calibri")
    txt(s, via, 4.6, ry, 3.4, 0.5, size=10.5, color=WHITE, vmid=True)
    txt(s, beh, 8.0, ry, 1.5, 0.5, size=9.5, color=RGBColor(0xB0,0xC8,0xE8), vmid=True, italic=True, font="Calibri")

txt(s, "Unique Constraints (UK)", 0.9, 3.35, 5, 0.3, size=13, bold=True, color=ORANGE, font="Calibri")
uks = [
    "Part.PartNo — รหัสอะไหล่ห้ามซ้ำ (DB index)",
    "AtmModelPart (AtmModelId, PartNo) — กันใส่อะไหล่ซ้ำในกลุ่ม",
    "EquivalentGroupMember (GroupId, PartNo) — กันสมาชิกซ้ำ",
    "Category.Name / Location.Code / Vendor.Code — ตรวจที่ controller",
]
for j, u in enumerate(uks):
    txt(s, f"•  {u}", 0.95, 3.72+j*0.34, 8.7, 0.32, size=10.5, color=WHITE, font="Calibri")

txt(s, "หมายเหตุ: ขนาด varchar เป็น design spec สำหรับ MySQL — ปัจจุบันรันบน SQLite (TEXT ไม่จำกัดความยาว)",
    0.9, 5.18, 8.8, 0.3, size=9, color=MUTED, italic=True, font="Calibri")

out = "D:/ATMApi/ATM-Inventory-System/ATM_Inventory_Phase1_ERD_Database.pptx"
prs.save(out)
print(f"Saved: {out}")
