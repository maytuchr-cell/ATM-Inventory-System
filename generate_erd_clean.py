"""
ATM Inventory — Clean ERD (compact boxes, tidy straight lines, clear arrowheads).
Lines never cross a box; relationship cardinality shown by markers (dot = 1, arrow = N).
Each box lists its KEY fields with a short description; full fields are in the Data Dictionary deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

NAVY=RGBColor(0x1C,0x35,0x57); BLUE=RGBColor(0x25,0x63,0xEB); TEAL=RGBColor(0x0D,0x94,0x88)
AMBER=RGBColor(0xB4,0x53,0x09); PURPLE=RGBColor(0x53,0x4A,0xB7); GREY=RGBColor(0x5A,0x67,0x78)
GREEN=RGBColor(0x05,0x96,0x69); WHITE=RGBColor(0xFF,0xFF,0xFF); TEXT=RGBColor(0x1E,0x2D,0x40)
MUTED=RGBColor(0x6B,0x72,0x80); LIGHT=RGBColor(0xF5,0xF7,0xFA); ORANGE=RGBColor(0xF5,0xA6,0x23)
ORANGED=RGBColor(0xD8,0x5A,0x30); GREY_L=RGBColor(0xD6,0xDA,0xE0); ROW_ALT=RGBColor(0xF6,0xF8,0xFB)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def slide():
    s=prs.slides.add_slide(prs.slide_layouts[6])
    f=s.background.fill; f.solid(); f.fore_color.rgb=LIGHT; return s
def box(s,x,y,w,h,fill,line=None,lw=0.75,radius=False):
    sh=s.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh
def txt(s,t,x,y,w,h,size=9,bold=False,color=TEXT,align=PP_ALIGN.LEFT,italic=False,vmid=False,font="Calibri"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.02);tf.margin_right=Inches(0.02);tf.margin_top=Inches(0.0);tf.margin_bottom=Inches(0.0)
    if vmid: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb
def header(s,title,sub):
    box(s,0,0,13.333,0.62,NAVY); box(s,0,0.62,13.333,0.03,ORANGE)
    txt(s,title,0.3,0.06,9.5,0.5,size=17,bold=True,color=WHITE,vmid=True)
    txt(s,sub,0.32,0.66,12.5,0.22,size=9,color=MUTED,italic=True)

ENT={}
RH=0.225
def entity(s,key,x,y,w,clr,bgc,fields,title=None):
    h=0.33+len(fields)*RH
    box(s,x,y,w,h,bgc,clr,1.25,radius=True)
    box(s,x,y,w,0.31,clr)
    txt(s,title or key,x+0.1,y+0.02,w-0.15,0.28,size=10.5,bold=True,color=WHITE,vmid=True)
    yy=y+0.34
    for i,(ic,fn,desc) in enumerate(fields):
        if i%2==1: box(s,x+0.03,yy,w-0.06,RH,ROW_ALT)
        if ic:
            kc={"PK":ORANGED,"UK":PURPLE,"FK":GREEN}.get(ic,MUTED)
            txt(s,ic,x+0.08,yy,0.32,RH,size=7,bold=True,color=kc,vmid=True)
        txt(s,fn,x+0.44,yy,w*0.40,RH,size=8.5,bold=(ic=="PK"),color=TEXT,vmid=True,font="Consolas")
        txt(s,desc,x+0.44+w*0.40,yy,w*0.60-0.48,RH,size=8,color=MUTED,vmid=True)
        yy+=RH
    ENT[key]=(x,y,w,h); return (x,y,w,h)

def _arrows(cn,clr):
    cn.line.color.rgb=clr; cn.line.width=Pt(1.75); cn.shadow.inherit=False
    ln=cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:headEnd'),{'type':'oval','w':'med','len':'med'}))
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'lg','len':'lg'}))

def link(s,a,b,clr=GREY):
    """Straight connector from a (the '1') to b (the 'N'), edge-to-edge on facing sides."""
    ax,ay,aw,ah=ENT[a]; bx,by,bw,bh=ENT[b]
    acx,acy=ax+aw/2,ay+ah/2; bcx,bcy=bx+bw/2,by+bh/2
    if abs(acx-bcx)>=abs(acy-bcy):
        x1,y1=(ax+aw,acy) if bcx>acx else (ax,acy)
        x2,y2=(bx,bcy) if bcx>acx else (bx+bw,bcy)
    else:
        x1,y1=(acx,ay+ah) if bcy>acy else (acx,ay)
        x2,y2=(bcx,by) if bcy>acy else (bcx,by+bh)
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    _arrows(cn,clr)

def comb(s,parent,children,busx,clr=GREY):
    """Tidy 'comb': one horizontal from parent.right, a vertical bus at busx, then a
       horizontal stub (with arrow) into each child's left edge. Children sit to the right."""
    px,py,pw,ph=ENT[parent]; pcy=py+ph/2
    ys=[ENT[c][1]+ENT[c][3]/2 for c in children]
    # horizontal from parent to bus (plain line, no arrow)
    seg=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(px+pw),Inches(pcy),Inches(busx),Inches(pcy))
    seg.line.color.rgb=clr; seg.line.width=Pt(1.75); seg.shadow.inherit=False
    # dot at parent end
    ln=seg.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:headEnd'),{'type':'oval','w':'med','len':'med'}))
    # vertical bus
    vb=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(busx),Inches(min(min(ys),pcy)),Inches(busx),Inches(max(max(ys),pcy)))
    vb.line.color.rgb=clr; vb.line.width=Pt(1.75); vb.shadow.inherit=False
    # stub to each child with arrow head
    for c,cy in zip(children,ys):
        cx=ENT[c][0]
        st=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(busx),Inches(cy),Inches(cx),Inches(cy))
        st.line.color.rgb=clr; st.line.width=Pt(1.75); st.shadow.inherit=False
        l2=st.line._get_or_add_ln(); l2.append(l2.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'lg','len':'lg'}))

def legend(s):
    y=7.18; x=0.4
    for tag,clr,d in [("PK",ORANGED,"Primary Key"),("UK",PURPLE,"Unique"),("FK",GREEN,"Foreign Key")]:
        txt(s,tag,x,y-0.02,0.35,0.22,size=8.5,bold=True,color=clr,vmid=True); x+=0.32
        txt(s,d,x,y-0.02,len(d)*0.072+0.3,0.22,size=8.5,color=MUTED,vmid=True); x+=len(d)*0.072+0.45
    txt(s,"ความสัมพันธ์ 1 : N   ●────▶   (วงกลม = ฝั่ง 1 / หัวลูกศร = ฝั่ง N)",x,y-0.02,5.5,0.22,size=8.5,color=GREY,italic=True,vmid=True)

# ══════════════════ SLIDE 1 — MASTER DATA ══════════════════
s=slide(); header(s,"Database Design — Master Data","ตารางข้อมูลหลัก + field สำคัญพร้อมคำอธิบาย   ·   field เต็มดูใน Data Dictionary")
# Part hub on the left-center; things that reference it placed to the right (comb)
entity(s,"Part",0.4,2.35,3.5,BLUE,RGBColor(0xE3,0xEC,0xFD),[
 ("PK","Id","รหัส (auto)"),("UK","PartNo","รหัสอะไหล่ (business key)"),("","PartName","ชื่อ/คำอธิบาย"),
 ("FK","CategoryId","→ Category (Sub Unit)"),("","Unit / MainUnit","หน่วย / กลุ่มหลัก"),
 ("","Min/Max/Reorder","นโยบายสต็อก"),("","IsActive / RowVersion","สถานะ / กันแก้ทับ")],title="Part  (ตัวอะไหล่)")
# Category is the PARENT of Part — placed left, links into Part
entity(s,"Category",0.4,0.95,3.5,PURPLE,RGBColor(0xF0,0xEF,0xFA),[
 ("PK","Id","รหัส"),("UK","Name","ชื่อหมวด (Sub Unit)"),("","IsActive","สถานะ")])
# Junction tables that reference Part — right side, stacked
entity(s,"AtmModel",9.5,0.95,3.45,PURPLE,RGBColor(0xF0,0xEF,0xFA),[
 ("PK","Id","รหัส"),("","ModelCode","รหัสกลุ่ม (Group Code)"),("","ModelName","ชื่อกลุ่ม/รุ่น"),("","Manufacturer","ผู้ผลิต")],title="AtmModel (ATM Group)")
entity(s,"AtmModelPart",5.7,2.5,3.2,TEAL,RGBColor(0xE9,0xF6,0xF4),[
 ("PK","Id","รหัส"),("FK","AtmModelId","→ AtmModel"),("FK","PartId","→ Part"),("","PartNo","snapshot")],title="AtmModelPart")
entity(s,"EquivalentGroup",9.5,4.55,3.45,PURPLE,RGBColor(0xF0,0xEF,0xFA),[
 ("PK","Id","รหัส"),("","Name","ชื่อกลุ่มอะไหล่ทดแทน"),("","CreatedAt","เวลาสร้าง")],title="EquivalentGroup")
entity(s,"EquivGroupMember",5.7,4.55,3.2,TEAL,RGBColor(0xE9,0xF6,0xF4),[
 ("PK","Id","รหัส"),("FK","GroupId","→ EquivalentGroup"),("FK","PartId","→ Part"),("","PartNo","snapshot")],title="EquivGroupMember")
# Standalone master tables (used by transactions, no Part link here)
entity(s,"Vendor",0.4,5.0,3.5,GREY_L if False else RGBColor(0xEE,0xF1,0xF6),GREY,[
 ("PK","Id","รหัส"),("UK","Code","รหัสผู้ขาย"),("","VendorType","GRG/LOCAL")],title="Vendor  (ใช้ใน Goods Receipt)")
# Links: Category(1)→Part(N) ; AtmModelPart(N)→Part(1) ; EquivGroupMember(N)→Part(1) ; AtmModel(1)→AtmModelPart(N) ; EquivalentGroup(1)→EquivGroupMember(N)
link(s,"Category","Part")            # vertical, Category above Part
link(s,"Part","AtmModelPart")        # Part.right → AtmModelPart.left  (Part is the 1 here)
link(s,"AtmModel","AtmModelPart")    # AtmModel above-right → AtmModelPart
link(s,"Part","EquivGroupMember")    # Part.right(lower) → EquivGroupMember.left
link(s,"EquivalentGroup","EquivGroupMember")
legend(s)

# ══════════════════ SLIDE 2 — STOCK & LEDGER ══════════════════
ENT.clear()
s=slide(); header(s,"Database Design — Stock & Ledger","สต็อกต่อคลัง · ชิ้นซีเรียล · บัญชีเคลื่อนไหว (Ledger)")
entity(s,"Part",0.4,2.9,2.9,BLUE,RGBColor(0xE3,0xEC,0xFD),[
 ("PK","Id","รหัส"),("UK","PartNo","รหัสอะไหล่"),("","PartName","ชื่อ"),("","(ยอด)","= SUM(PartStock)")],title="Part")
# children of Part, stacked right, left-aligned for a clean comb
entity(s,"PartStock",5.0,1.0,4.2,TEAL,RGBColor(0xE9,0xF6,0xF4),[
 ("PK","Id","รหัส"),("FK","PartId","→ Part"),("FK","LocationId","→ Location"),
 ("","GoodQty / DefectiveQty","ดี / เสีย"),("","UpdatedAt / RowVersion","เวลายอด / กันทับ · UQ(Part,Loc)")],title="PartStock  (ยอดในคลัง)")
entity(s,"PartUnit",5.0,3.15,4.2,AMBER,RGBColor(0xFB,0xF1,0xE7),[
 ("PK","Id","รหัส"),("FK","PartId","→ Part"),("FK","LocationId","→ Location"),
 ("UK","SerialNo","ซีเรียล (ห้ามซ้ำ)"),("","Condition / Status","สภาพ / InStock·Issued·Disposed")],title="PartUnit  (ชิ้นซีเรียล)")
entity(s,"StockMovement",5.0,5.25,4.2,TEAL,RGBColor(0xE9,0xF6,0xF4),[
 ("PK","Id","รหัส"),("FK","PartId","→ Part (Restrict)"),("FK","PartUnitId","→ PartUnit (ถ้ามีซีเรียล)"),
 ("","Type / Qty / Cond","ประเภท/จำนวน/สภาพ"),("","From/ToLocationId","คลังออก/เข้า · User · Time")],title="StockMovement  (Ledger)")
entity(s,"Location",10.0,3.15,2.9,BLUE,RGBColor(0xE3,0xEC,0xFD),[
 ("PK","Id","รหัสคลัง"),("UK","Code","รหัส"),("","Name","ชื่อคลัง")],title="Location")
comb(s,"Part",["PartStock","PartUnit","StockMovement"],4.4)   # Part → 3 children via tidy bus
link(s,"PartUnit","StockMovement")                            # PartUnit(1) → StockMovement(N)
txt(s,"Location เชื่อม PartStock / PartUnit / StockMovement ผ่านฟิลด์ LocationId (ดูป้าย → Location ในกล่อง)",
    10.0,6.3,3.0,0.6,size=8,color=MUTED,italic=True)
legend(s)

# ══════════════════ SLIDE 3 — TRANSACTIONS ══════════════════
ENT.clear()
s=slide(); header(s,"Database Design — Transactions","รับเข้า · เบิก/คืน · โอน · นับสต็อก · ทำลาย   ·   ทุกตารางอ้าง Part ด้วย FK PartId")
# Three independent parent→child pairs, laid out in 3 columns (no crossing)
entity(s,"GoodsReceipt",0.4,1.0,3.9,AMBER,RGBColor(0xFB,0xF1,0xE7),[
 ("PK","Id","รหัส"),("","ReceiptNo","เลขใบรับ"),("","Source","GRG/LocalVendor"),
 ("FK","VendorId","→ Vendor"),("FK","LocationId","→ Location"),("","HandlingCost","ค่าจัดการ")],title="GoodsReceipt  (ใบรับเข้า)")
entity(s,"GoodsReceiptLine",0.4,4.0,3.9,AMBER,RGBColor(0xFB,0xF1,0xE7),[
 ("PK","Id","รหัส"),("FK","GoodsReceiptId","→ GoodsReceipt"),("FK","PartId","→ Part"),
 ("","Qty / Condition","จำนวน / สภาพ"),("","SerialNo","ซีเรียล → สร้าง PartUnit")],title="GoodsReceiptLine")

entity(s,"Ticket",4.7,1.0,3.9,GREEN,RGBColor(0xEA,0xF6,0xF0),[
 ("PK","TicketId","รหัส"),("","TechUser...","ข้อมูลช่าง"),("","Req/ApprovedPartNo","อะไหล่ขอ/อนุมัติ"),
 ("","Status","Pending/Approved/.."),("","AttachmentPath","รูปอาการเสีย")],title="Ticket  (ใบเบิก)")
entity(s,"ReturnRequest",4.7,4.0,3.9,GREEN,RGBColor(0xEA,0xF6,0xF0),[
 ("PK","Id","รหัส"),("FK","TicketId","→ Ticket (บังคับ)"),("FK","PartId","→ Part"),
 ("","Condition","สภาพที่คืน"),("","Location From/To","ต้นทาง/ปลายทาง")],title="ReturnRequest  (คืน)")

entity(s,"StockCount",9.0,1.0,3.9,AMBER,RGBColor(0xFB,0xF1,0xE7),[
 ("PK","Id","รหัส"),("","CountType","Cycle/Annual"),("","Period","งวด"),("","Status","Draft/../Completed")],title="StockCount  (นับสต็อก)")
entity(s,"StockCountLine",9.0,4.0,3.9,AMBER,RGBColor(0xFB,0xF1,0xE7),[
 ("PK","Id","รหัส"),("FK","StockCountId","→ StockCount"),("FK","PartId","→ Part"),
 ("","System/PhysicalQty","ยอดระบบ/นับจริง"),("","Variance","ผลต่าง")],title="StockCountLine")

link(s,"GoodsReceipt","GoodsReceiptLine")
link(s,"Ticket","ReturnRequest")
link(s,"StockCount","StockCountLine")
legend(s)
txt(s,"StockTransfer และ DisposalRequest (ไม่แสดงในผังนี้) ก็อ้าง Part ด้วย FK PartId เช่นกัน — ดู Data Dictionary",
    0.4,6.65,12.5,0.22,size=8.5,color=MUTED,italic=True)

out="D:/ATMApi/ATM-Inventory-System/ATM_Inventory_ERD_Clean.pptx"
prs.save(out); print("Saved:",out,"slides:",len(prs.slides._sldIdLst))
