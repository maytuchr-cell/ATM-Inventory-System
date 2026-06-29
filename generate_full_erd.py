"""
ATM Inventory System — FULL SYSTEM ERD (all 23 tables)
Solid line = enforced FK · Dashed line = loose link by PartNo string (no DB FK)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR

NAVY=RGBColor(0x1C,0x35,0x57); BLUE=RGBColor(0x25,0x63,0xEB); TEAL=RGBColor(0x0D,0x94,0x88)
AMBER=RGBColor(0xB4,0x53,0x09); PURPLE=RGBColor(0x53,0x4A,0xB7); GREY=RGBColor(0x64,0x74,0x8B)
GREEN=RGBColor(0x05,0x96,0x69); WHITE=RGBColor(0xFF,0xFF,0xFF); TEXT=RGBColor(0x1E,0x2D,0x40)
MUTED=RGBColor(0x6B,0x72,0x80); LIGHT=RGBColor(0xF5,0xF7,0xFA); ORANGE=RGBColor(0xF5,0xA6,0x23)
GREY_L=RGBColor(0xD1,0xD5,0xDB)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6])
bg=s.background.fill; bg.solid(); bg.fore_color.rgb=LIGHT

def box(x,y,w,h,fill,line=None,lw=0.75,radius=False):
    sh=s.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def txt(t,x,y,w,h,size=9,bold=False,color=TEXT,align=PP_ALIGN.LEFT,italic=False,vmid=False,font="Calibri"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.02);tf.margin_right=Inches(0.02);tf.margin_top=Inches(0.0);tf.margin_bottom=Inches(0.0)
    if vmid: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb

# header
box(0,0,13.333,0.62,NAVY); box(0,0.62,13.333,0.03,ORANGE)
txt("ATM Inventory System — Full Database ERD (23 tables)",0.3,0.06,9,0.5,size=18,bold=True,color=WHITE,vmid=True)
txt("เส้นทึบ = FK บังคับจริง   ·   เส้นประ = เชื่อมหลวมด้วย PartNo (string, ไม่มี FK)",0.3,0.4,9,0.2,size=9,color=RGBColor(0xB0,0xC8,0xE8),italic=True)

ENT={}  # key -> (cx_top, x,y,w,h)
def entity(key,x,y,w,clr,bgc,fields,title=None):
    rh=0.17; h=0.30+len(fields)*rh
    box(x,y,w,h,bgc,clr,1.0,radius=True)
    box(x,y,w,0.28,clr)
    txt(title or key,x+0.07,y+0.01,w-0.1,0.26,size=8.5,bold=True,color=WHITE,vmid=True)
    yy=y+0.31
    for ic,fn in fields:
        if ic:
            kc={"PK":ORANGE,"UK":PURPLE,"FK":GREEN}.get(ic,MUTED)
            txt(ic,x+0.05,yy,0.3,rh,size=6.5,bold=True,color=kc,vmid=True)
        txt(fn,x+(0.34 if ic else 0.09),yy,w-0.4,rh,size=7,color=TEXT,vmid=True,font="Consolas")
        yy+=rh
    ENT[key]=(x,y,w,h)

# cluster labels
def clabel(t,x,y,clr): txt(t,x,y,3,0.2,size=8.5,bold=True,color=clr)

# ===== MASTER DATA (blue) — top-left =====
clabel("MASTER DATA",0.3,0.74,BLUE)
entity("Category",0.3,0.98,1.45,BLUE,RGBColor(0xEE,0xF3,0xFE),[("PK","Id"),("UK","Name"),("","IsActive")])
entity("Vendor",1.9,0.98,1.45,BLUE,RGBColor(0xEE,0xF3,0xFE),[("PK","Id"),("UK","Code"),("","Name/Type")])
entity("Location",3.5,0.98,1.45,BLUE,RGBColor(0xEE,0xF3,0xFE),[("PK","Id"),("UK","Code"),("","Type")])
entity("AtmModel",5.1,0.98,1.45,BLUE,RGBColor(0xEE,0xF3,0xFE),[("PK","Id"),("","ModelCode"),("","ModelName")])

# ===== SYSTEM (grey) — top-right =====
clabel("SYSTEM",11.0,0.74,GREY)
entity("User",11.0,0.98,2.0,GREY,RGBColor(0xF1,0xF2,0xF4),[("PK","Id"),("UK","Email"),("","Role")])
entity("AuditLog",11.0,1.85,2.0,GREY,RGBColor(0xF1,0xF2,0xF4),[("PK","Id"),("","EntityType/Id"),("","Action")])
entity("SystemSettings",11.0,2.72,2.0,GREY,RGBColor(0xF1,0xF2,0xF4),[("PK","Id"),("","IsFrozen")])

# ===== PART hub (blue, emphasized) =====
entity("Part",3.4,2.35,1.85,BLUE,RGBColor(0xE3,0xEC,0xFD),
       [("PK","Id"),("UK","PartNo"),("FK","CategoryId"),("","PartName"),("","min/max/reorder")],title="Part  (ตัวอะไหล่)")

# ===== STOCK (teal) — right of Part =====
clabel("STOCK / INVENTORY",6.9,2.12,TEAL)
entity("PartStock",6.9,2.35,1.7,TEAL,RGBColor(0xE9,0xF6,0xF4),[("PK","Id"),("FK","PartId"),("FK","LocationId"),("","Good/DefQty")])
entity("PartUnit",8.75,2.35,1.7,TEAL,RGBColor(0xE9,0xF6,0xF4),[("PK","Id"),("FK","PartId"),("FK","LocationId"),("UK","SerialNo")])
entity("StockMovement",6.9,3.55,1.7,TEAL,RGBColor(0xE9,0xF6,0xF4),[("PK","Id"),("FK","PartId  *NEW*"),("","From/ToLoc"),("","Type/Qty")])

# ===== CONFIG (purple) — left-mid =====
clabel("CONFIG",0.3,2.12,PURPLE)
entity("AtmModelPart",5.1,2.35,1.55,PURPLE,RGBColor(0xF0,0xEF,0xFA),[("PK","Id"),("FK","AtmModelId"),("~","PartNo")])
entity("EquivalentGroup",0.3,2.35,1.55,PURPLE,RGBColor(0xF0,0xEF,0xFA),[("PK","Id"),("","Name")])
entity("EquivGroupMember",0.3,3.25,1.55,PURPLE,RGBColor(0xF0,0xEF,0xFA),[("PK","Id"),("FK","GroupId"),("~","PartNo")])
entity("EquivalentPart",1.95,3.25,1.45,PURPLE,RGBColor(0xF0,0xEF,0xFA),[("PK","Id"),("~","Orig/Equiv")])

# ===== TRANSACTIONS (amber) — bottom =====
clabel("TRANSACTIONS",0.3,4.62,AMBER)
entity("GoodsReceipt",0.3,4.85,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("FK","VendorId"),("FK","LocationId")])
entity("GoodsReceiptLine",2.15,4.85,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("FK","GoodsReceiptId"),("~","PartNo")])
entity("Ticket",0.3,6.0,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","TicketId"),("~","Req/AppPartNo")])
entity("ReturnRequest",2.15,6.0,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("FK","TicketId"),("~","PartNo")])
entity("StockTransfer",4.0,4.85,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("~","PartNo"),("~","From/ToLoc")])
entity("DisposalRequest",4.0,6.0,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("~","PartNo"),("~","LocationId")])
entity("StockCount",5.85,4.85,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("","Type/Status")])
entity("StockCountLine",5.85,6.0,1.7,AMBER,RGBColor(0xFB,0xF1,0xE7),[("PK","Id"),("FK","StockCountId"),("~","PartNo")])

def C(k):  # center of an entity
    x,y,w,h=ENT[k]; return (x+w/2,y+h/2)
def edge(a,b,clr,dashed=False,lw=1.25):
    (x1,y1),(x2,y2)=C(a),C(b)
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=clr; cn.line.width=Pt(lw); cn.shadow.inherit=False
    if dashed:
        from pptx.oxml.ns import qn
        ln=cn.line._get_or_add_ln(); d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)

# ---- SOLID FK edges ----
for a,b in [("Category","Part"),("Part","PartStock"),("Part","PartUnit"),("Part","StockMovement"),
            ("Location","PartStock"),("AtmModel","AtmModelPart"),
            ("EquivalentGroup","EquivGroupMember"),("GoodsReceipt","GoodsReceiptLine"),
            ("Vendor","GoodsReceipt"),("Ticket","ReturnRequest"),("StockCount","StockCountLine")]:
    edge(a,b,GREY,dashed=False,lw=1.5)

# ---- DASHED loose (PartNo) edges to Part ----
for a in ["AtmModelPart","EquivGroupMember","GoodsReceiptLine","ReturnRequest","StockTransfer",
          "DisposalRequest","StockCountLine","Ticket"]:
    edge(a,"Part",RGBColor(0xC0,0x6A,0x2A),dashed=True,lw=1.0)

# legend
ly=7.12
box(0.3,ly,0.5,0.04,GREY); txt("FK (enforced)",0.85,ly-0.08,1.6,0.2,size=8,color=MUTED,vmid=True)
box(2.6,ly,0.5,0.04,RGBColor(0xC0,0x6A,0x2A)); txt("loose link by PartNo (no FK)",3.15,ly-0.08,3,0.2,size=8,color=MUTED,vmid=True)
txt("~ = อ้างด้วย string PartNo / int LocationId",6.4,ly-0.08,4,0.2,size=8,color=MUTED,italic=True,vmid=True)

# ════════════════════════════════════════════════════════════════════════
#  DATA DICTIONARY SLIDES  (Field | Type | Null | Key | Description)
# ════════════════════════════════════════════════════════════════════════
ROW_ALT=RGBColor(0xF3,0xF6,0xFB); ORANGED=RGBColor(0xD8,0x5A,0x30)

def blank():
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    f=sl.background.fill; f.solid(); f.fore_color.rgb=LIGHT
    return sl

def dbox(sl,x,y,w,h,fill,line=None,lw=0.75,radius=False):
    sh=sl.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh

def dtxt(sl,t,x,y,w,h,size=9,bold=False,color=TEXT,align=PP_ALIGN.LEFT,italic=False,vmid=False,font="Calibri"):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.03);tf.margin_right=Inches(0.03);tf.margin_top=Inches(0.0);tf.margin_bottom=Inches(0.0)
    if vmid: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=t
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; r.font.name=font
    return tb

def dheader(sl,title,subtitle=None):
    dbox(sl,0,0,13.333,0.62,NAVY); dbox(sl,0,0.62,13.333,0.03,ORANGE)
    dtxt(sl,title,0.3,0.06,12.5,0.5,size=17,bold=True,color=WHITE,vmid=True)
    if subtitle: dtxt(sl,subtitle,0.32,0.66,12,0.22,size=9,color=MUTED,italic=True)

def ftable(sl,x,y,w,name,clr,fields,row_h=0.235,fsize=8.5,note=None):
    fw=[w*0.23,w*0.19,w*0.07,w*0.07,w*0.44]
    dtxt(sl,name,x,y-0.26,w,0.24,size=11,bold=True,color=clr)
    dbox(sl,x,y,w,0.28,clr)
    cx=x
    for i,cn in enumerate(["Field","Data Type","Null","Key","Description"]):
        al=PP_ALIGN.CENTER if cn in("Null","Key") else PP_ALIGN.LEFT
        dtxt(sl,cn,cx+0.04,y+0.01,fw[i]-0.06,0.26,size=8,bold=True,color=WHITE,vmid=True,align=al)
        cx+=fw[i]
    dbox(sl,x,y+0.28,w,len(fields)*row_h,WHITE,GREY_L,0.75)
    for ri,(fn,ft,nl,ky,ds) in enumerate(fields):
        ry=y+0.28+ri*row_h
        if ri%2==1: dbox(sl,x+0.02,ry,w-0.04,row_h,ROW_ALT)
        cx=x
        dtxt(sl,fn,cx+0.05,ry,fw[0]-0.08,row_h,size=fsize,bold=True,color=TEXT,vmid=True,font="Consolas"); cx+=fw[0]
        dtxt(sl,ft,cx+0.05,ry,fw[1]-0.08,row_h,size=fsize-0.5,color=BLUE,vmid=True,font="Consolas"); cx+=fw[1]
        nc=MUTED if nl else RGBColor(0xB0,0xB0,0xB0)
        dtxt(sl,"Y" if nl else "—",cx,ry,fw[2],row_h,size=7,color=nc,vmid=True,align=PP_ALIGN.CENTER); cx+=fw[2]
        kc={"PK":ORANGED,"UK":PURPLE,"FK":GREEN,"~":AMBER}.get(ky,MUTED)
        if ky: dtxt(sl,ky,cx,ry,fw[3],row_h,size=fsize-0.5,bold=True,color=kc,vmid=True,align=PP_ALIGN.CENTER); cx+=fw[3]
        else: cx+=fw[3]
        dtxt(sl,ds,cx+0.05,ry,fw[4]-0.08,row_h,size=fsize,color=TEXT,vmid=True)
    if note:
        ny=y+0.28+len(fields)*row_h+0.05
        dtxt(sl,note,x,ny,w,0.3,size=8,color=TEAL,italic=True)

def dlegend(sl):
    items=[("PK",ORANGED,"Primary Key"),("UK",PURPLE,"Unique"),("FK",GREEN,"Foreign Key"),("~",AMBER,"loose link by PartNo (no FK)")]
    x=0.4; y=7.18
    for tag,clr,desc in items:
        dtxt(sl,tag,x,y-0.02,0.3,0.22,size=8.5,bold=True,color=clr,vmid=True); x+=0.3
        dtxt(sl,desc,x,y-0.02,len(desc)*0.07+0.4,0.22,size=8.5,color=MUTED,vmid=True); x+=len(desc)*0.07+0.55

# ---- field definitions (field, type(MySQL), nullable, key, desc-TH) ----
T={
"Part":(BLUE,[
 ("Id","int",False,"PK","Primary key"),("OrderNumber","varchar(50)",False,"","เลขสั่งซื้อ/อ้างอิงเอกสาร"),
 ("PartNo","varchar(50)",False,"UK","รหัสอะไหล่ — business key (ห้ามซ้ำ)"),("PartName","varchar(255)",False,"","ชื่อ/คำอธิบายอะไหล่"),
 ("Unit","varchar(20)",False,"","หน่วยนับ (default pcs)"),("SerialNo","varchar(100)",True,"","ซีเรียล (legacy → ใช้ PartUnit)"),
 ("CategoryId","int",True,"FK","→ Category (Sub Unit)"),("CatalogueRef","varchar(100)",True,"","อ้างอิงเลขใน catalogue"),
 ("MinStock","int",False,"","สต็อกขั้นต่ำ"),("MaxStock","int",False,"","สต็อกขั้นสูง"),("ReorderPoint","int",False,"","จุดสั่งซื้อใหม่"),
 ("TrackingNumber","varchar(100)",True,"","เลขติดตามขนส่ง (legacy)"),("Aging","int",True,"","อายุสต็อก วัน (legacy/คำนวณ)"),
 ("CostPerUnit","decimal(18,2)",True,"","ต้นทุนต่อหน่วย"),("IsActive","tinyint(1)",False,"","สถานะ — false = soft delete"),
 ("CreatedAt","datetime",False,"","เวลาสร้าง"),("UpdatedAt","datetime",False,"","เวลาแก้ไขล่าสุด"),
 ("ExpiryDate","datetime",True,"","วันหมดอายุ (legacy → PartUnit)"),("IsUnrepairable","tinyint(1)",False,"","ซ่อมไม่ได้ (legacy → PartUnit)"),
 ("MainUnit","varchar(100)",True,"","กลุ่มหลัก เช่น Cabinet"),("Remark","varchar(500)",True,"","หมายเหตุ"),
 ("ImagePath","varchar(255)",True,"","พาธรูป /assets/parts/..."),
 ]),
"Category":(BLUE,[("Id","int",False,"PK","Primary key"),("Name","varchar(100)",False,"UK","ชื่อหมวด = Sub Unit (ห้ามซ้ำ)"),
 ("Description","varchar(500)",True,"","รายละเอียด"),("IsActive","tinyint(1)",False,"","สถานะ")]),
"Vendor":(BLUE,[("Id","int",False,"PK","Primary key"),("Name","varchar(150)",False,"","ชื่อผู้ขาย"),
 ("Code","varchar(50)",False,"UK","รหัสผู้ขาย (ห้ามซ้ำ)"),("VendorType","varchar(20)",False,"","GRG | LOCAL"),
 ("ContactInfo","varchar(255)",True,"","ข้อมูลติดต่อ"),("IsActive","tinyint(1)",False,"","สถานะ")]),
"Location":(BLUE,[("Id","int",False,"PK","Primary key"),("Name","varchar(150)",False,"","ชื่อคลัง/จุดเก็บ"),
 ("Code","varchar(50)",False,"UK","รหัสคลัง (ห้ามซ้ำ)"),("LocationType","varchar(30)",False,"","ประเภท เช่น WH/GRG/SCRAP"),
 ("IsActive","tinyint(1)",False,"","สถานะ")]),
"AtmModel":(PURPLE,[("Id","int",False,"PK","Primary key"),("ModelCode","varchar(50)",False,"","รหัสกลุ่ม/รุ่น (Group Code)"),
 ("ModelName","varchar(150)",False,"","ชื่อกลุ่ม/รุ่น"),("Manufacturer","varchar(100)",True,"","ผู้ผลิต เช่น GRG/NCR"),
 ("Description","varchar(500)",True,"","รายละเอียด"),("IsActive","tinyint(1)",False,"","สถานะ")]),
"AtmModelPart":(PURPLE,[("Id","int",False,"PK","Primary key"),("AtmModelId","int",False,"FK","→ AtmModel (Cascade)"),
 ("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo — UK(AtmModelId,PartNo)")]),
"EquivalentGroup":(PURPLE,[("Id","int",False,"PK","Primary key"),("Name","varchar(150)",False,"","ชื่อกลุ่มอะไหล่ทดแทน"),
 ("Description","varchar(500)",True,"","รายละเอียด"),("CreatedAt","datetime",False,"","เวลาสร้าง")]),
"EquivalentGroupMember":(PURPLE,[("Id","int",False,"PK","Primary key"),("GroupId","int",False,"FK","→ EquivalentGroup (Cascade)"),
 ("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo — UK(GroupId,PartNo)")]),
"EquivalentPart":(PURPLE,[("Id","int",False,"PK","Primary key"),("OriginalPartNo","varchar(50)",False,"~","อะไหล่ต้นทาง"),
 ("EquivalentPartNo","varchar(50)",False,"~","อะไหล่ทดแทน — UK(Orig,Equiv)")]),
"PartStock":(TEAL,[("Id","int",False,"PK","Primary key"),("PartId","int",False,"FK","→ Part"),
 ("LocationId","int",False,"FK","→ Location"),("GoodQty","int",False,"","จำนวนสภาพดี"),
 ("DefectiveQty","int",False,"","จำนวนสภาพเสีย")]),
"PartUnit":(TEAL,[("Id","int",False,"PK","Primary key"),("PartId","int",False,"FK","→ Part (Cascade)"),
 ("LocationId","int",True,"FK","→ Location (SetNull)"),("SerialNo","varchar(100)",False,"UK","ซีเรียลของชิ้นนี้ (ห้ามซ้ำ)"),
 ("Condition","varchar(20)",False,"","Good | Defective"),("ExpiryDate","datetime",True,"","วันหมดอายุของชิ้นนี้"),
 ("IsUnrepairable","tinyint(1)",False,"","ซ่อมไม่ได้ → Disposal"),("ReceivedAt","datetime",False,"","วันรับเข้า (ใช้คำนวณ Aging)"),
 ("Status","varchar(20)",False,"","InStock|Issued|Disposed")]),
"StockMovement":(TEAL,[("Id","int",False,"PK","Primary key"),("MovementType","varchar(20)",False,"","GR/Issue/Return/Transfer/Disposal/Adjust"),
 ("PartId","int",False,"FK","→ Part (Restrict) — NEW"),("PartNo","varchar(50)",False,"","snapshot ของ PartNo ตอนนั้น"),
 ("FromLocationId","int",True,"","คลังต้นทาง (ออก)"),("ToLocationId","int",True,"","คลังปลายทาง (เข้า)"),
 ("Qty","int",False,"","จำนวนที่ขยับ"),("Condition","varchar(20)",False,"","Good | Defective"),
 ("RefType","varchar(20)",True,"","Ticket/GoodsReceipt/Transfer..."),("RefId","varchar(50)",True,"","เลขเอกสารอ้างอิง"),
 ("Cost","decimal(18,2)",True,"","ต้นทุน ณ ตอนขยับ"),("SerialNo","varchar(100)",True,"","ซีเรียล (ถ้ามี)"),
 ("Remarks","varchar(500)",True,"","หมายเหตุ"),("UserName","varchar(100)",False,"","ผู้ทำรายการ"),
 ("Timestamp","datetime",False,"","เวลาที่ทำรายการ")]),
"GoodsReceipt":(AMBER,[("Id","int",False,"PK","Primary key"),("ReceiptNo","varchar(50)",False,"","เลขที่ใบรับ"),
 ("Source","varchar(20)",False,"","GRG | LocalVendor"),("VendorId","int",True,"FK","→ Vendor"),
 ("RefDocument","varchar(100)",True,"","Forecast/Lot/PO"),("LocationId","int",False,"FK","→ Location ปลายทาง"),
 ("ReceivedBy","varchar(100)",False,"","ผู้รับ"),("ReceivedAt","datetime",False,"","เวลารับ"),
 ("HandlingCost","decimal(18,2)",True,"","ค่าจัดการ/ขนส่ง")]),
"GoodsReceiptLine":(AMBER,[("Id","int",False,"PK","Primary key"),("GoodsReceiptId","int",False,"FK","→ GoodsReceipt (Cascade)"),
 ("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo"),("Qty","int",False,"","จำนวน"),
 ("Condition","varchar(20)",False,"","Good | Defective"),("SerialNo","varchar(100)",True,"","ซีเรียล"),
 ("IsManualAdjust","tinyint(1)",False,"","ปรับมือ (ต้องมี Remarks)"),("Remarks","varchar(500)",True,"","หมายเหตุ")]),
"Ticket":(AMBER,[("TicketId","int",False,"PK","Primary key"),("TechEmail","varchar(150)",False,"","อีเมลช่าง"),
 ("TechId","varchar(50)",False,"","รหัสช่าง"),("TechName","varchar(150)",False,"","ชื่อช่าง"),
 ("TechPhone","varchar(50)",False,"","เบอร์ช่าง"),("TechDept","varchar(100)",False,"","แผนก"),
 ("RequestedPartNo","varchar(50)",True,"~","อะไหล่ที่ขอ (อ้าง Part.PartNo)"),("ApprovedPartNo","varchar(50)",True,"~","อะไหล่ที่อนุมัติ"),
 ("FaultySerialNo","varchar(100)",True,"","ซีเรียลของเสีย"),("FaultyPartNo","varchar(50)",True,"","รหัสอะไหล่เสีย"),
 ("MachineModel","varchar(100)",True,"","รุ่นเครื่อง ATM"),("Description","varchar(500)",True,"","อาการ/รายละเอียด"),
 ("AttachmentPath","varchar(255)",True,"","รูปแนบ"),("MainCause","varchar(255)",True,"","สาเหตุหลัก (FR-IW-01)"),
 ("LogisticsCost","decimal(18,2)",True,"","ค่าขนส่งจริง"),("Status","varchar(20)",False,"","Pending/Approved/..."),
 ("IsDOA","tinyint(1)",False,"","Dead-on-arrival (FR-MC-04)"),("CreatedAt","datetime",False,"","เวลาสร้าง"),
 ("ReceivedAt","datetime",True,"","เวลารับของ"),("DueDate","datetime",True,"","กำหนดส่ง")]),
"ReturnRequest":(AMBER,[("Id","int",False,"PK","Primary key"),("TicketId","int",False,"FK","→ Ticket (Restrict, บังคับ)"),
 ("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo"),("Condition","varchar(20)",False,"","Good | Defective"),
 ("SourceType","varchar(20)",False,"","Technician/GRG/LocalVendor"),("LocationFromId","int",False,"","คลังต้นทาง"),
 ("LocationToId","int",False,"","คลังปลายทาง"),("ReturnedBy","varchar(100)",False,"","ผู้คืน"),
 ("CreatedAt","datetime",False,"","เวลาสร้าง")]),
"StockTransfer":(AMBER,[("Id","int",False,"PK","Primary key"),("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo"),
 ("Qty","int",False,"","จำนวน"),("Condition","varchar(20)",False,"","Good | Defective"),
 ("FromLocationId","int",False,"~","คลังต้นทาง"),("ToLocationId","int",False,"~","คลังปลายทาง"),
 ("Status","varchar(20)",False,"","Pending/Approved/InTransit/Received"),("RequestedBy","varchar(100)",False,"","ผู้ขอ"),
 ("ApprovedBy","varchar(100)",True,"","ผู้อนุมัติ"),("CreatedAt","datetime",False,"","เวลาสร้าง"),
 ("ApprovedAt","datetime",True,"","เวลาอนุมัติ"),("ConfirmedAt","datetime",True,"","เวลายืนยันส่ง"),
 ("ReceivedAt","datetime",True,"","เวลารับปลายทาง")]),
"StockCount":(AMBER,[("Id","int",False,"PK","Primary key"),("CountType","varchar(20)",False,"","Cycle | Annual"),
 ("Period","varchar(20)",False,"","เช่น 2026-Q1"),("Status","varchar(20)",False,"","Draft/InProgress/Completed"),
 ("IsSystemFrozen","tinyint(1)",False,"","freeze ระบบระหว่างนับ"),("StartedBy","varchar(100)",False,"","ผู้เริ่มนับ"),
 ("CreatedAt","datetime",False,"","เวลาสร้าง"),("CompletedAt","datetime",True,"","เวลาเสร็จ")]),
"StockCountLine":(AMBER,[("Id","int",False,"PK","Primary key"),("StockCountId","int",False,"FK","→ StockCount (Cascade)"),
 ("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo"),("LocationId","int",False,"","คลังที่นับ"),
 ("SystemQty","int",False,"","ยอดระบบ ณ ตอนเริ่มนับ"),("PhysicalQty","int",True,"","ยอดนับจริง"),
 ("Variance","int (computed)",False,"","ผลต่าง = Physical - System (ไม่เก็บ)"),("AdjustApproved","tinyint(1)",False,"","อนุมัติปรับยอด"),
 ("Remarks","varchar(500)",True,"","หมายเหตุ")]),
"DisposalRequest":(AMBER,[("Id","int",False,"PK","Primary key"),("PartNo","varchar(50)",False,"~","อ้าง Part.PartNo"),
 ("SerialNo","varchar(100)",True,"","ซีเรียล"),("LocationId","int",False,"~","คลัง (ปกติ Scrap)"),
 ("Qty","int",False,"","จำนวน"),("Status","varchar(20)",False,"","Pending/Approved/Disposed"),
 ("ReasonCode","varchar(30)",False,"","Expired/Unrepairable/Damaged/Other"),("RequestedBy","varchar(100)",False,"","ผู้ขอ"),
 ("ApprovedBy","varchar(100)",True,"","ผู้อนุมัติ"),("CreatedAt","datetime",False,"","เวลาสร้าง"),
 ("ApprovedAt","datetime",True,"","เวลาอนุมัติ"),("DisposedAt","datetime",True,"","เวลาทำลาย")]),
"User":(GREY,[("Id","int",False,"PK","Primary key"),("Email","varchar(150)",False,"UK","อีเมล (lower-case, ห้ามซ้ำ)"),
 ("PasswordHash","varchar(255)",False,"","salt:hash (PBKDF2)"),("Role","varchar(20)",False,"","Admin | Tech"),
 ("Name","varchar(150)",False,"","ชื่อผู้ใช้"),("IsActive","tinyint(1)",False,"","สถานะ"),("CreatedAt","datetime",False,"","เวลาสร้าง")]),
"AuditLog":(GREY,[("Id","int",False,"PK","Primary key"),("EntityType","varchar(30)",False,"","Part/Category/Location/..."),
 ("EntityId","varchar(50)",False,"","id ของ entity"),("Action","varchar(20)",False,"","CREATE/UPDATE/DELETE/APPROVE"),
 ("OldValues","text",True,"","ค่าก่อนแก้ (JSON)"),("NewValues","text",True,"","ค่าหลังแก้ (JSON)"),
 ("UserId","varchar(50)",False,"","id ผู้ทำ"),("UserName","varchar(100)",False,"","ชื่อผู้ทำ"),("Timestamp","datetime",False,"","เวลา")]),
"SystemSettings":(GREY,[("Id","int",False,"PK","Primary key (row เดียว Id=1)"),("IsFrozen","tinyint(1)",False,"","freeze ระบบ (ระหว่างนับสต็อก)"),
 ("ActiveStockCountId","int",True,"","รอบนับที่กำลัง active")]),
}

# ---- slide layouts: list of (title, subtitle, [ (tableName, x, y, w), ... ]) ----
slides=[
 ("Data Dictionary — Part (ตัวอะไหล่)","ตาราง Parts · ยอดคงเหลือไม่เก็บที่นี่ — คำนวณจาก SUM(PartStock.GoodQty)",
   [("Part",0.6,1.15,12.1)]),
 ("Data Dictionary — Master Data (Category · Vendor · Location)",None,
   [("Category",0.4,1.2,4.1),("Vendor",4.65,1.2,4.1),("Location",8.9,1.2,4.0)]),
 ("Data Dictionary — Part Configuration (ATM Group · Equivalent)",None,
   [("AtmModel",0.4,1.2,6.1),("AtmModelPart",0.4,4.0,6.1),
    ("EquivalentGroup",6.8,1.2,6.1),("EquivalentGroupMember",6.8,3.4,6.1),("EquivalentPart",6.8,5.2,6.1)]),
 ("Data Dictionary — Stock Core (PartStock · PartUnit)",None,
   [("PartStock",0.5,1.3,6.0),("PartUnit",6.8,1.3,6.0)]),
 ("Data Dictionary — StockMovement (Ledger)","FK PartId เชื่อม Part จริง (Restrict) · PartNo เป็น snapshot",
   [("StockMovement",0.7,1.3,11.9)]),
 ("Data Dictionary — Goods Receipt",None,
   [("GoodsReceipt",0.5,1.3,6.0),("GoodsReceiptLine",6.8,1.3,6.0)]),
 ("Data Dictionary — Ticket (Issue & Withdrawal)",None,
   [("Ticket",0.7,1.2,11.9)]),
 ("Data Dictionary — Returns · Transfers",None,
   [("ReturnRequest",0.5,1.3,6.0),("StockTransfer",6.8,1.3,6.0)]),
 ("Data Dictionary — Stock Count · Disposal",None,
   [("StockCount",0.4,1.3,4.1),("StockCountLine",4.65,1.3,4.2),("DisposalRequest",9.0,1.3,4.0)]),
 ("Data Dictionary — System (User · AuditLog · Settings)",None,
   [("User",0.4,1.3,4.1),("AuditLog",4.65,1.3,4.2),("SystemSettings",9.0,1.3,4.0)]),
]
for title,sub,tabs in slides:
    sl=blank(); dheader(sl,title,sub)
    for nm,x,y,w in tabs:
        clr,flds=T[nm]
        ftable(sl,x,y,w,nm,clr,flds)
    dlegend(sl)

out="D:/ATMApi/ATM-Inventory-System/ATM_Inventory_FullSystem_ERD.pptx"
prs.save(out); print("Saved:",out,"slides:",len(prs.slides._sldIdLst))
