"""ATM Inventory — Infrastructure / Deployment Diagram (IIS + SSL + API + MySQL + Assets)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

NAVY=RGBColor(0x1C,0x35,0x57); BLUE=RGBColor(0x25,0x63,0xEB); TEAL=RGBColor(0x0D,0x94,0x88)
AMBER=RGBColor(0xB4,0x53,0x09); PURPLE=RGBColor(0x53,0x4A,0xB7); GREY=RGBColor(0x5A,0x67,0x78)
GREEN=RGBColor(0x05,0x96,0x69); WHITE=RGBColor(0xFF,0xFF,0xFF); TEXT=RGBColor(0x1E,0x2D,0x40)
MUTED=RGBColor(0x6B,0x72,0x80); LIGHT=RGBColor(0xF5,0xF7,0xFA); ORANGE=RGBColor(0xF5,0xA6,0x23)
SKY=RGBColor(0xE3,0xEC,0xFD); MINT=RGBColor(0xE9,0xF6,0xF4); SAND=RGBColor(0xFB,0xF1,0xE7)
GREYBG=RGBColor(0xEE,0xF1,0xF6); GREEN_L=RGBColor(0xEA,0xF6,0xF0)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6]); f=s.background.fill; f.solid(); f.fore_color.rgb=LIGHT

def box(x,y,w,h,fill,line=None,lw=1.0,radius=True):
    sh=s.shapes.add_shape(5 if radius else 1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh
def txt(t,x,y,w,h,size=9,bold=False,color=TEXT,align=PP_ALIGN.LEFT,italic=False,vmid=False,font="Calibri"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.04);tf.margin_right=Inches(0.04);tf.margin_top=Inches(0.01);tf.margin_bottom=Inches(0.01)
    if vmid: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,line in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=line; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name=font
    return tb
def node(x,y,w,h,title,lines,clr,bgc):
    box(x,y,w,h,bgc,clr,1.25)
    box(x,y,w,0.30,clr)
    txt(title,x+0.1,y+0.01,w-0.15,0.28,size=10,bold=True,color=WHITE,vmid=True)
    txt(lines,x+0.12,y+0.36,w-0.2,h-0.42,size=8.5,color=TEXT)
def arrow(x1,y1,x2,y2,clr=GREY,lw=2.0,dash=False,label=None,lx=None,ly=None):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=clr; cn.line.width=Pt(lw); cn.shadow.inherit=False
    ln=cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'lg','len':'lg'}))
    if dash: ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    if label: txt(label,lx,ly,2.2,0.24,size=8,bold=True,color=clr,align=PP_ALIGN.CENTER,italic=True)

# header
box(0,0,13.333,0.62,NAVY,radius=False); box(0,0.62,13.333,0.03,ORANGE,radius=False)
txt("ATM Inventory — Infrastructure / Deployment",0.3,0.06,10,0.5,size=18,bold=True,color=WHITE,vmid=True)
txt("Web (IIS) + Mobile → HTTPS → ASP.NET Core API → MySQL + Asset storage",0.32,0.66,12,0.22,size=9.5,color=MUTED,italic=True)

# ── CLIENTS (left) ──
txt("CLIENTS",0.4,1.05,3,0.24,size=9,bold=True,color=MUTED)
node(0.4,1.35,3.0,1.35,"🖥  Web Browser",
     "Admin console\n• SystemAdmin / Staff / Auditor\n• admin*.html (static)",BLUE,SKY)
node(0.4,3.15,3.0,1.35,"📱  Mobile App",
     "Technician (role=Tech)\n• สร้าง/ติดตาม Ticket\n• ถ่ายรูป defect",TEAL,MINT)
txt("ผู้ใช้ในองค์กร (Intranet)",0.4,4.62,3,0.22,size=8,color=MUTED,italic=True)

# ── HTTPS boundary ──
box(4.05,1.2,0.9,3.6,RGBColor(0xFF,0xF4,0xE0),ORANGE,1.0)
txt("🔒\nHTTPS\nSSL cert\n(JWT\nBearer)",4.05,1.5,0.9,3.0,size=8.5,bold=True,color=AMBER,align=PP_ALIGN.CENTER,vmid=True)

# ── SERVER zone (IIS) ──
box(5.25,1.05,5.0,4.35,RGBColor(0xEC,0xEF,0xF4),NAVY,1.25)
txt("WINDOWS SERVER  ·  IIS",5.4,1.12,4.7,0.24,size=9.5,bold=True,color=NAVY)
node(5.5,1.5,4.5,0.95,"IIS  (Web Server + Reverse Proxy)",
     "• เสิร์ฟ static frontend (HTML/JS/CSS)\n• ASP.NET Core Module → host API (in-process)",GREY,GREYBG)
node(5.5,2.75,4.5,1.15,"ASP.NET Core API  (.NET 10, Kestrel)",
     "• Controllers (Auth/Parts/Ticket/GoodsReceipt/...)\n• JWT auth + Role policies (SystemAdmin/Staff/Auditor/Tech)\n• EF Core  ·  StockService · AuditService",BLUE,SKY)
node(5.5,4.15,4.5,1.1,"Serves  /assets/*  (static files)",
     "รูปอะไหล่ → อ่านจากโฟลเดอร์ภายนอก (AssetPath)\nconfig: appsettings.Production.json",AMBER,SAND)

# ── DATA zone (right) ──
txt("DATA (บน server / โฟลเดอร์แยก)",10.6,1.05,2.7,0.24,size=9,bold=True,color=MUTED)
node(10.6,1.5,2.55,1.5,"🗄  Database",
     "MySQL (production)\nDev: SQLite (AtmInventory.db)\n• EF Core / migrations\n• FK + concurrency (RowVersion)",TEAL,MINT)
node(10.6,3.25,2.55,1.35,"🖼  Asset Storage",
     "D:\\ATMAssets\\parts\\\n• รูปอะไหล่ (นอก repo)\n• IIS/API serve /assets",GREY,GREYBG)

# ── arrows ──
arrow(3.4,2.0,4.05,2.4,GREEN,label="fetch (web)",lx=2.9,ly=1.55)
arrow(3.4,3.8,4.05,3.6,TEAL,label="REST + JWT (mobile)",lx=2.7,ly=4.02)
arrow(4.95,3.0,5.5,3.2,ORANGE)                 # HTTPS → API
arrow(10.0,3.3,10.6,2.4,BLUE)   # API → DB
txt("EF Core",9.35,2.15,1.2,0.22,size=8,bold=True,color=BLUE,align=PP_ALIGN.CENTER,italic=True)
arrow(10.0,4.7,10.6,4.0,AMBER)  # assets(API) → storage
txt("อ่านรูป",9.35,4.7,1.2,0.22,size=8,bold=True,color=AMBER,align=PP_ALIGN.CENTER,italic=True)

# footer notes
txt("Deploy checklist:  appsettings.Production.json (MySQL conn + AssetPath)  ·  web.config (ASP.NET Core Module)  ·  SSL binding  ·  CORS จำกัด origin  ·  EF Migrations แทน EnsureCreated",
    0.4,6.05,12.7,0.5,size=9,color=NAVY,italic=True)
txt("Legend:  ── เส้นทึบ = การเรียกจริง (HTTP/EF)   ·   ทุกการเรียก API ผ่าน HTTPS + JWT",
    0.4,6.7,12.7,0.3,size=8.5,color=MUTED,italic=True)

out="D:/ATMApi/ATM-Inventory-System/ATM_Inventory_Infrastructure.pptx"
prs.save(out); print("Saved:",out)
