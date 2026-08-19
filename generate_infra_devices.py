"""ATM Inventory — Physical Infrastructure / Network Topology (device-style)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

NAVY=RGBColor(0x1C,0x35,0x57); BLUE=RGBColor(0x25,0x63,0xEB); TEAL=RGBColor(0x0D,0x94,0x88)
AMBER=RGBColor(0xB4,0x53,0x09); GREY=RGBColor(0x5A,0x67,0x78); GREEN=RGBColor(0x18,0x9E,0x5B)
WHITE=RGBColor(0xFF,0xFF,0xFF); TEXT=RGBColor(0x1E,0x2D,0x40); MUTED=RGBColor(0x6B,0x72,0x80)
LIGHT=RGBColor(0xF5,0xF7,0xFA); ORANGE=RGBColor(0xF5,0xA6,0x23); RED=RGBColor(0xC0,0x39,0x2B)
DARK=RGBColor(0x22,0x2A,0x38); STEEL=RGBColor(0xC7,0xCF,0xDA); PANEL=RGBColor(0xE7,0xEC,0xF2)
SCREEN=RGBColor(0x1E,0x40,0x66); LED=RGBColor(0x35,0xD0,0x7A)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6]); f=s.background.fill; f.solid(); f.fore_color.rgb=LIGHT

def rect(x,y,w,h,fill,line=None,lw=1.0,radius=False,shape=None):
    sh=s.shapes.add_shape(shape if shape else (5 if radius else 1),Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False; return sh
def txt(t,x,y,w,h,size=9,bold=False,color=TEXT,align=PP_ALIGN.CENTER,italic=False,vmid=True,font="Calibri"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.02);tf.margin_right=Inches(0.02);tf.margin_top=Inches(0.0);tf.margin_bottom=Inches(0.0)
    if vmid: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for i,ln in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=ln; r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
        r.font.color.rgb=color; r.font.name=font
    return tb
def line(x1,y1,x2,y2,clr=GREY,lw=2.0,dash=False):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=clr; cn.line.width=Pt(lw); cn.shadow.inherit=False
    if dash:
        from pptx.oxml.ns import qn
        ln=cn.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    return cn

# ── device icons ─────────────────────────────────────────────────────────
def monitor(cx,y,label,sub,clr=BLUE):
    w,h=1.15,0.78
    rect(cx-w/2,y,w,h,DARK,clr,1.5,radius=True)         # bezel
    rect(cx-w/2+0.07,y+0.07,w-0.14,h-0.14,SCREEN)       # screen
    rect(cx-0.11,y+h,0.22,0.12,STEEL)                    # stand
    rect(cx-0.28,y+h+0.12,0.56,0.07,GREY)                # base
    txt(label,cx-0.9,y+h+0.22,1.8,0.24,size=9,bold=True,color=TEXT)
    txt(sub,cx-0.9,y+h+0.44,1.8,0.34,size=7.5,color=MUTED)
def phone(cx,y,label,sub,clr=TEAL):
    w,h=0.5,0.92
    rect(cx-w/2,y,w,h,DARK,clr,1.5,radius=True)
    rect(cx-w/2+0.05,y+0.1,w-0.1,h-0.2,SCREEN)
    rect(cx-0.06,y+0.03,0.12,0.03,GREY)                  # speaker
    txt(label,cx-0.9,y+h+0.05,1.8,0.24,size=9,bold=True,color=TEXT)
    txt(sub,cx-0.9,y+h+0.27,1.8,0.34,size=7.5,color=MUTED)
def server(cx,y,label,sub,clr=NAVY,h=1.5):
    w=1.1
    rect(cx-w/2,y,w,h,PANEL,clr,1.5,radius=True)
    for i in range(4):                                    # rack slots
        sy=y+0.14+i*(h-0.2)/4
        rect(cx-w/2+0.1,sy,w-0.2,(h-0.3)/4*0.72,RGBColor(0xD3,0xDA,0xE3),GREY,0.5)
        rect(cx+w/2-0.26,sy+0.02,0.05,0.05,LED)           # LED
    txt(label,cx-0.95,y+h+0.05,1.9,0.24,size=9,bold=True,color=TEXT)
    txt(sub,cx-0.95,y+h+0.27,1.9,0.5,size=7.5,color=MUTED)
def database(cx,y,label,sub,clr=TEAL):
    w,h=1.05,1.25
    rect(cx-w/2,y,w,h,RGBColor(0xE9,0xF6,0xF4),clr,1.75,shape=MSO_SHAPE.CAN)
    txt("🛢",cx-0.4,y+0.25,0.8,0.6,size=22,color=clr)
    txt(label,cx-0.95,y+h+0.02,1.9,0.24,size=9,bold=True,color=TEXT)
    txt(sub,cx-0.95,y+h+0.24,1.9,0.5,size=7.5,color=MUTED)
def storage(cx,y,label,sub):
    w,h=1.1,1.15
    rect(cx-w/2,y,w,h,PANEL,GREY,1.5,radius=True)
    for i in range(3):
        rect(cx-w/2+0.12,y+0.16+i*0.34,w-0.24,0.24,RGBColor(0xD3,0xDA,0xE3),GREY,0.5)
        rect(cx+w/2-0.28,y+0.22+i*0.34,0.06,0.06,LED)
    txt(label,cx-0.95,y+h+0.02,1.9,0.24,size=9,bold=True,color=TEXT)
    txt(sub,cx-0.95,y+h+0.24,1.9,0.5,size=7.5,color=MUTED)
def netbox(cx,y,w,h,label,sub,clr,emoji):
    rect(cx-w/2,y,w,h,RGBColor(0xFF,0xF4,0xE0) if clr==ORANGE else PANEL,clr,1.75,radius=True)
    txt(emoji,cx-w/2,y+0.05,w,0.34,size=15)
    txt(label,cx-w/2,y+0.40,w,0.26,size=9,bold=True,color=clr)
    txt(sub,cx-w/2,y+h-0.24,w,0.22,size=7.5,color=MUTED)

# ── header ──
rect(0,0,13.333,0.62,NAVY); rect(0,0.62,13.333,0.03,ORANGE)
txt("ATM Inventory — Infrastructure (Physical / Network Topology)",0.3,0.06,10.5,0.5,size=17,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
txt("อุปกรณ์จริง: เครื่องผู้ใช้ → เครือข่ายองค์กร → ห้อง Server (IIS · MySQL · Storage)",0.32,0.66,12.5,0.22,size=9,color=MUTED,align=PP_ALIGN.LEFT,italic=True)

# ── zone backgrounds ──
rect(0.25,1.05,3.3,5.6,RGBColor(0xEE,0xF2,0xF7),STEEL,1.0,radius=True)
txt("สำนักงาน / หน้างาน  (Client Devices)",0.35,1.15,3.1,0.3,size=9.5,bold=True,color=NAVY,align=PP_ALIGN.LEFT)
rect(4.0,1.05,3.0,5.6,RGBColor(0xF1,0xEE,0xF7),RGBColor(0xB9,0xB0,0xD8),1.0,radius=True)
txt("เครือข่ายองค์กร (LAN / Intranet)",4.1,1.15,2.8,0.3,size=9.5,bold=True,color=RGBColor(0x53,0x4A,0xB7),align=PP_ALIGN.LEFT)
rect(7.4,1.05,5.65,5.6,RGBColor(0xEC,0xF4,0xF1),RGBColor(0x9C,0xC7,0xBB),1.0,radius=True)
txt("ห้อง Server / Data Center",7.5,1.15,5.4,0.3,size=9.5,bold=True,color=TEAL,align=PP_ALIGN.LEFT)

# ── CLIENT devices ──
monitor(1.9,1.7,"Admin PC ×N","SystemAdmin / Staff / Auditor\nWeb browser")
phone(1.35,4.15,"Mobile","ช่าง (Tech)")
phone(2.45,4.15,"Tablet","ช่าง (Tech)")
txt("ช่างสร้าง/ติดตาม Ticket + ถ่ายรูป",0.4,5.9,3.0,0.4,size=7.5,color=MUTED,italic=True)

# ── NETWORK ──
netbox(5.5,1.75,2.2,0.95,"Firewall + SSL","HTTPS · JWT Bearer",ORANGE,"🔥")
netbox(5.5,3.5,2.2,0.9,"Switch / Router","LAN",GREY,"🔀")
txt("Wi-Fi (มือถือ)",4.15,4.6,2.0,0.3,size=8,color=MUTED,italic=True)

# ── SERVERS ──
server(8.5,1.7,"Web / App Server","IIS + ASP.NET Core API\n(.NET 10, in-process)",NAVY,1.6)
database(11.4,1.75,"Database Server","MySQL (production)\nFK · concurrency")
storage(11.4,4.2,"File Storage","D:\\ATMAssets\\parts\n(รูปอะไหล่ · /assets)")
server(8.5,4.3,"— (เดียวกับ Web)","API serve /assets +\nคุยกับ DB",NAVY,1.35) if False else None
txt("Web + API รันบนเครื่องเดียว (IIS host in-process)",7.5,6.25,5.5,0.3,size=8,color=MUTED,italic=True,align=PP_ALIGN.LEFT)

# ── connections ──
line(1.9,2.7,5.5,2.2,GREEN,2.25)                 # Admin PC → firewall
line(1.9,4.6,4.5,3.7,TEAL,2.0,dash=True)         # phones (wifi) → switch
line(2.45,4.6,4.6,3.8,TEAL,2.0,dash=True)
line(5.5,2.65,5.5,3.5,ORANGE,2.25)               # firewall → switch
line(6.6,3.9,8.5,2.5,GREY,2.25)                  # switch → web server
line(9.05,2.4,10.9,2.3,BLUE,2.25)                # web server → DB
line(9.05,3.0,10.9,4.6,AMBER,2.25)               # web server → storage
# small cardinal labels
txt("HTTPS",3.3,2.05,1.4,0.24,size=7.5,bold=True,color=GREEN,italic=True)
txt("EF Core",9.5,2.0,1.3,0.24,size=7.5,bold=True,color=BLUE,italic=True)
txt("อ่านรูป",9.5,3.65,1.3,0.24,size=7.5,bold=True,color=AMBER,italic=True)

# ── footer ──
txt("Deploy: Windows Server + IIS (ASP.NET Core Module) · SSL cert (มีแล้ว) · MySQL · appsettings.Production.json (conn + AssetPath) · CORS จำกัด origin",
    0.3,6.85,12.9,0.35,size=8.5,color=NAVY,italic=True,align=PP_ALIGN.LEFT)

out="D:/ATMApi/ATM-Inventory-System/ATM_Inventory_Infrastructure_Devices.pptx"
prs.save(out); print("Saved:",out)
